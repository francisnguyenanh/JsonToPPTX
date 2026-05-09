# tests/test_shapes.py
import pytest
from pptx import Presentation
from pptx.util import Emu
from renderer import geometry as G
from renderer.shapes import (
    add_gradient_fill, add_textbox_styled, add_badge,
    add_dot_bullet, add_separator_line, add_semantic_icon,
    _set_gradient_fill_on_spPr, _set_solid_fill_on_spPr
)
from pptx.oxml.ns import qn


def blank_slide():
    prs = Presentation()
    prs.slide_width = Emu(G.SLIDE_W)
    prs.slide_height = Emu(G.SLIDE_H)
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def test_add_textbox_styled_returns_shape():
    slide = blank_slide()
    tb = add_textbox_styled(slide, G.pt(10), G.pt(10), G.pt(200), G.pt(30),
                            "Hello", size_pt=12, color_hex="172759")
    assert tb is not None
    assert tb.text_frame.paragraphs[0].runs[0].text == "Hello"


def test_add_textbox_bodypr_is_explicit():
    slide = blank_slide()
    tb = add_textbox_styled(slide, G.pt(10), G.pt(10), G.pt(200), G.pt(30),
                            "Test", inset=G.INS_CARD, autofit="norm")
    bodyPr = tb.text_frame._txBody.find(qn("a:bodyPr"))
    assert bodyPr is not None
    assert bodyPr.get("wrap") is not None
    assert bodyPr.get("anchor") is not None
    # Check autofit child present
    has_autofit = (
        bodyPr.find(qn("a:normAutofit")) is not None or
        bodyPr.find(qn("a:noAutofit")) is not None or
        bodyPr.find(qn("a:spAutoFit")) is not None
    )
    assert has_autofit


def test_add_badge_creates_shape():
    slide = blank_slide()
    badge = add_badge(slide, G.pt(10), G.pt(10), "1", "2362B0", size_pt=14)
    assert badge is not None
    assert badge.text_frame.paragraphs[0].runs[0].text == "1"


def test_add_dot_bullet_creates_shape():
    slide = blank_slide()
    dot = add_dot_bullet(slide, G.pt(10), G.pt(10), "4A9EE0", 6)
    assert dot is not None


def test_add_separator_line():
    slide = blank_slide()
    line = add_separator_line(slide, G.pt(20), G.pt(100), G.pt(200), G.pt(1), "B8D0EE")
    assert line is not None


def test_gradient_fill_on_spPr():
    slide = blank_slide()
    shape = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(G.pt(100)), Emu(G.pt(50)))
    spPr = shape._element.spPr
    _set_gradient_fill_on_spPr(spPr, "FFFFFF", "EEF4FB", 9000000)
    gradFill = spPr.find(qn("a:gradFill"))
    assert gradFill is not None
    gsLst = gradFill.find(qn("a:gsLst"))
    assert len(gsLst) == 2
    lin = gradFill.find(qn("a:lin"))
    assert lin.get("ang") == "9000000"


def test_semantic_icon_returns_shapes():
    slide = blank_slide()
    for itype in ("ai", "settings", "data", "doc", "team", "check", "default"):
        shapes = add_semantic_icon(slide, G.pt(10), G.pt(10), 36, "4A9EE0", itype)
        assert len(shapes) >= 1, f"icon_type={itype} returned no shapes"
