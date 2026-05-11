"""svg_parser.py — Convert SVG strings into python-pptx slides.

Supported SVG elements
----------------------
<rect>            → rectangle / rounded-rect shape with solid or gradient fill
<circle>          → oval shape
<ellipse>         → oval shape
<line>            → straight connector
<text> / <tspan>  → textbox (with run-level bold, italic, color, size)
<image>           → picture (http/https URL or data URI)
<g>               → group — applies cumulative translate transform to children
<linearGradient>  → gradient fill, referenced via fill="url(#id)"

Coordinate system
-----------------
viewBox="0 0 W H" is parsed and scaled so W×H maps to SLIDE_W×SLIDE_H.
If no viewBox, 1 SVG unit = 1pt = 12700 EMU (assumes 1280×720 canvas).

Public API
----------
parse_svg_to_slide(prs, svg_string) -> None
    Appends one slide to the given Presentation, populated from the SVG.
"""
from __future__ import annotations

import io
import math
import re
import urllib.request
from typing import Optional
from xml.etree import ElementTree as ET

from lxml import etree as lxml_etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from . import geometry as G
from .shapes import (
    _set_gradient_fill_on_spPr,
    _set_solid_fill_on_spPr,
    _hex_upper,
)

# ── Namespaces ────────────────────────────────────────────────────────────────
_SVG = "http://www.w3.org/2000/svg"
_XLINK = "http://www.w3.org/1999/xlink"
ET.register_namespace("", _SVG)
ET.register_namespace("xlink", _XLINK)

_SLIDE_W_PT = 1280  # pts
_SLIDE_H_PT = 720
_PT_TO_EMU = G.PT   # 12700

# ── OOXML gradient angles (1/60000 degree units, from design_system.py) ──────
_ANG_TOP_BOTTOM = 9000000
_ANG_BTM_TOP    = 16200000
_ANG_LEFT_RIGHT = 5400000
_ANG_VERTICAL   = 0


# ─────────────────────────────────────────────────────────────────────────────
# Colour helpers
# ─────────────────────────────────────────────────────────────────────────────

_NAMED_COLORS: dict[str, Optional[str]] = {
    "white": "FFFFFF", "black": "000000", "red": "FF0000",
    "green": "008000", "blue": "0000FF", "gray": "808080",
    "grey": "808080", "orange": "FFA500", "yellow": "FFFF00",
    "purple": "800080", "violet": "EE82EE", "cyan": "00FFFF",
    "magenta": "FF00FF", "pink": "FFC0CB", "brown": "A52A2A",
    "silver": "C0C0C0", "gold": "FFD700", "navy": "000080",
    "teal": "008080", "lime": "00FF00", "indigo": "4B0082",
    "transparent": None, "none": None,
}


def _parse_color(val: str) -> Optional[str]:
    """Return uppercase 6-char hex string, or None for none/transparent."""
    if not val:
        return None
    val = val.strip()
    if val in ("none", "transparent"):
        return None
    if val.startswith("url("):
        # Gradient reference — handled separately
        return None
    if val.startswith("#"):
        h = val[1:]
        if len(h) == 3:
            h = h[0] * 2 + h[1] * 2 + h[2] * 2
        return h.upper()[:6]
    m = re.match(r"rgb\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)", val)
    if m:
        return "%02X%02X%02X" % (int(m[1]), int(m[2]), int(m[3]))
    m = re.match(r"rgba\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*,\s*[\d.]+\s*\)", val)
    if m:
        return "%02X%02X%02X" % (int(m[1]), int(m[2]), int(m[3]))
    return _NAMED_COLORS.get(val.lower())


def _parse_style(style_str: str) -> dict[str, str]:
    result: dict[str, str] = {}
    for part in style_str.split(";"):
        part = part.strip()
        if ":" in part:
            k, v = part.split(":", 1)
            result[k.strip().lower()] = v.strip()
    return result


def _attr(el: ET.Element, *names: str, style: dict | None = None) -> str:
    """Read attribute from element, checking inline style first."""
    if style:
        for name in names:
            if name in style:
                return style[name]
    for name in names:
        v = el.get(name, "")
        if v:
            return v
    return ""


# ─────────────────────────────────────────────────────────────────────────────
# Gradient reference helper
# ─────────────────────────────────────────────────────────────────────────────

def _gradient_id(fill: str) -> Optional[str]:
    """Extract '#id' from 'url(#id)', or None."""
    m = re.match(r"url\(#([^)]+)\)", fill.strip())
    return m.group(1) if m else None


def _svg_grad_angle(x1: float, y1: float, x2: float, y2: float) -> int:
    """Map SVG linearGradient x1,y1→x2,y2 to OOXML ang (1/60000 degrees)."""
    dx = x2 - x1
    dy = y2 - y1
    if abs(dy) >= abs(dx):
        return _ANG_TOP_BOTTOM if dy >= 0 else _ANG_BTM_TOP
    else:
        return _ANG_LEFT_RIGHT if dx >= 0 else _ANG_VERTICAL


def _pct(val: str) -> float:
    """Parse '50%' → 0.5, or '0.5' → 0.5."""
    val = val.strip()
    if val.endswith("%"):
        return float(val[:-1]) / 100.0
    return float(val)


# ─────────────────────────────────────────────────────────────────────────────
# Transform parsing
# ─────────────────────────────────────────────────────────────────────────────

def _parse_transform(transform: str) -> tuple[float, float]:
    """Extract cumulative (tx, ty) translation from a transform string.
    Handles translate(), matrix(), and combinations."""
    tx, ty = 0.0, 0.0
    if not transform:
        return tx, ty
    # matrix(a,b,c,d,e,f) → e=tx, f=ty
    for m in re.finditer(r"matrix\(([^)]+)\)", transform):
        parts = re.split(r"[\s,]+", m.group(1).strip())
        if len(parts) == 6:
            tx += float(parts[4])
            ty += float(parts[5])
    # translate(x[,y])
    for m in re.finditer(r"translate\(([^)]+)\)", transform):
        parts = re.split(r"[\s,]+", m.group(1).strip())
        tx += float(parts[0])
        ty += float(parts[1]) if len(parts) > 1 else 0.0
    return tx, ty


# ─────────────────────────────────────────────────────────────────────────────
# Font size parsing
# ─────────────────────────────────────────────────────────────────────────────

def _parse_font_size(val: str, default: float = 16.0) -> float:
    """Parse '24px', '24pt', '24' → float pt."""
    val = val.strip()
    if not val:
        return default
    if val.endswith("px"):
        return float(val[:-2]) * 0.75  # 1px ≈ 0.75pt
    if val.endswith("pt"):
        return float(val[:-2])
    if val.endswith("em"):
        return float(val[:-2]) * default
    try:
        return float(val)
    except ValueError:
        return default


# ─────────────────────────────────────────────────────────────────────────────
# Shape fill helpers
# ─────────────────────────────────────────────────────────────────────────────

def _apply_fill(spPr, fill_val: str, gradients: dict, opacity: float = 1.0):
    """Apply solid or gradient fill to an spPr element."""
    gid = _gradient_id(fill_val)
    if gid and gid in gradients:
        g = gradients[gid]
        _set_gradient_fill_on_spPr(spPr, g["from"], g["to"], g["angle"])
        return

    color = _parse_color(fill_val)
    if color:
        alpha = int((1.0 - opacity) * 100000) if opacity < 1.0 else None
        _set_solid_fill_on_spPr(spPr, color, alpha)
    else:
        # no fill
        for tag in (qn("a:noFill"), qn("a:solidFill"), qn("a:gradFill")):
            for el in spPr.findall(tag):
                spPr.remove(el)
        lxml_etree.SubElement(spPr, qn("a:noFill"))


def _remove_line(spPr):
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        spPr.remove(ln)
    lxml_etree.SubElement(spPr, qn("a:ln")).append(
        lxml_etree.fromstring(
            '<a:noFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
        )
    )


def _set_stroke(spPr, stroke_hex: str, width_pt: float):
    """Apply a solid stroke to a shape."""
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        spPr.remove(ln)
    w_emu = int(width_pt * 12700)
    ln = lxml_etree.SubElement(spPr, qn("a:ln"), w=str(w_emu))
    solidFill = lxml_etree.SubElement(ln, qn("a:solidFill"))
    lxml_etree.SubElement(solidFill, qn("a:srgbClr"), val=_hex_upper(stroke_hex))


# ─────────────────────────────────────────────────────────────────────────────
# Coordinate conversion
# ─────────────────────────────────────────────────────────────────────────────

class _Ctx:
    """Parsing context: holds scale factors and gradient definitions."""

    def __init__(self, scale_x: float, scale_y: float, gradients: dict):
        self.scale_x = scale_x
        self.scale_y = scale_y
        self.gradients = gradients

    def emu_x(self, val: float) -> int:
        return int(val * self.scale_x * _PT_TO_EMU)

    def emu_y(self, val: float) -> int:
        return int(val * self.scale_y * _PT_TO_EMU)

    def emu_w(self, val: float) -> int:
        return max(1, int(val * self.scale_x * _PT_TO_EMU))

    def emu_h(self, val: float) -> int:
        return max(1, int(val * self.scale_y * _PT_TO_EMU))

    def pt_font(self, val: float) -> float:
        """Scale font size by average of x/y scale factors."""
        return val * (self.scale_x + self.scale_y) / 2


# ─────────────────────────────────────────────────────────────────────────────
# Element renderers
# ─────────────────────────────────────────────────────────────────────────────

def _render_rect(slide, el: ET.Element, ctx: _Ctx, tx: float, ty: float):
    sty = _parse_style(el.get("style", ""))
    x = float(el.get("x", 0)) + tx
    y = float(el.get("y", 0)) + ty
    w = float(el.get("width", 0))
    h = float(el.get("height", 0))
    if w <= 0 or h <= 0:
        return

    rx = el.get("rx") or el.get("ry") or sty.get("rx", "0")
    try:
        rx_val = float(rx)
    except (ValueError, TypeError):
        rx_val = 0.0

    fill_val   = sty.get("fill", el.get("fill", "none"))
    stroke_val = sty.get("stroke", el.get("stroke", "none"))
    stroke_w   = float(sty.get("stroke-width", el.get("stroke-width", "1")) or 1)
    opacity    = float(sty.get("opacity", el.get("opacity", "1")) or 1)

    shape = slide.shapes.add_shape(
        1,
        Emu(ctx.emu_x(x)), Emu(ctx.emu_y(y)),
        Emu(ctx.emu_w(w)), Emu(ctx.emu_h(h)),
    )
    spPr = shape._element.spPr

    # Rounded corners
    if rx_val > 0:
        pg = spPr.find(qn("a:prstGeom"))
        if pg is not None:
            pg.set("prst", "roundRect")
            av = pg.find(qn("a:avLst"))
            if av is None:
                av = lxml_etree.SubElement(pg, qn("a:avLst"))
            for gd in av.findall(qn("a:gd")):
                av.remove(gd)
            # adj = rx / min(w,h) * 2 * 100000, capped at 50000
            adj = min(50000, int(rx_val / min(w, h) * 2 * 100000))
            lxml_etree.SubElement(av, qn("a:gd"), name="adj", fmla=f"val {adj}")

    # Fill
    _apply_fill(spPr, fill_val, ctx.gradients, opacity)

    # Stroke
    stroke_color = _parse_color(stroke_val)
    if stroke_color:
        _set_stroke(spPr, stroke_color, stroke_w)
    else:
        _remove_line(spPr)

    shape.text_frame.text = ""


def _render_oval(slide, el: ET.Element, ctx: _Ctx, tx: float, ty: float):
    tag = el.tag.split("}")[-1]
    sty = _parse_style(el.get("style", ""))

    if tag == "circle":
        cx = float(el.get("cx", 0)) + tx
        cy = float(el.get("cy", 0)) + ty
        r  = float(el.get("r", 10))
        x, y, w, h = cx - r, cy - r, r * 2, r * 2
    else:  # ellipse
        cx = float(el.get("cx", 0)) + tx
        cy = float(el.get("cy", 0)) + ty
        rx = float(el.get("rx", 10))
        ry = float(el.get("ry", 10))
        x, y, w, h = cx - rx, cy - ry, rx * 2, ry * 2

    if w <= 0 or h <= 0:
        return

    fill_val   = sty.get("fill", el.get("fill", "none"))
    stroke_val = sty.get("stroke", el.get("stroke", "none"))
    stroke_w   = float(sty.get("stroke-width", el.get("stroke-width", "1")) or 1)
    opacity    = float(sty.get("opacity", el.get("opacity", "1")) or 1)

    shape = slide.shapes.add_shape(
        9,  # MSO oval
        Emu(ctx.emu_x(x)), Emu(ctx.emu_y(y)),
        Emu(ctx.emu_w(w)), Emu(ctx.emu_h(h)),
    )
    spPr = shape._element.spPr
    _apply_fill(spPr, fill_val, ctx.gradients, opacity)

    stroke_color = _parse_color(stroke_val)
    if stroke_color:
        _set_stroke(spPr, stroke_color, stroke_w)
    else:
        _remove_line(spPr)

    shape.text_frame.text = ""


def _render_line(slide, el: ET.Element, ctx: _Ctx, tx: float, ty: float):
    sty = _parse_style(el.get("style", ""))
    x1 = float(el.get("x1", 0)) + tx
    y1 = float(el.get("y1", 0)) + ty
    x2 = float(el.get("x2", 0)) + tx
    y2 = float(el.get("y2", 0)) + ty

    stroke_val = sty.get("stroke", el.get("stroke", "808080"))
    stroke_w   = float(sty.get("stroke-width", el.get("stroke-width", "1")) or 1)
    stroke_color = _parse_color(stroke_val) or "808080"

    # Represent as a thin rect
    if abs(y2 - y1) <= stroke_w * 2:
        # Horizontal line
        lx = min(x1, x2)
        shape = slide.shapes.add_shape(
            1,
            Emu(ctx.emu_x(lx)), Emu(ctx.emu_y(y1 - stroke_w / 2)),
            Emu(ctx.emu_w(abs(x2 - x1))), Emu(ctx.emu_h(max(1, stroke_w))),
        )
    elif abs(x2 - x1) <= stroke_w * 2:
        # Vertical line
        ly = min(y1, y2)
        shape = slide.shapes.add_shape(
            1,
            Emu(ctx.emu_x(x1 - stroke_w / 2)), Emu(ctx.emu_y(ly)),
            Emu(ctx.emu_w(max(1, stroke_w))), Emu(ctx.emu_h(abs(y2 - y1))),
        )
    else:
        # Diagonal: draw a thin rect from top-left of bounding box
        lx, ly = min(x1, x2), min(y1, y2)
        shape = slide.shapes.add_shape(
            1,
            Emu(ctx.emu_x(lx)), Emu(ctx.emu_y(ly)),
            Emu(ctx.emu_w(abs(x2 - x1))), Emu(max(1, ctx.emu_h(stroke_w))),
        )

    spPr = shape._element.spPr
    _set_solid_fill_on_spPr(spPr, stroke_color)
    _remove_line(spPr)
    shape.text_frame.text = ""


def _render_text(slide, el: ET.Element, ctx: _Ctx, tx: float, ty: float):
    """Render <text> element. Handles inline text and <tspan> children."""
    sty = _parse_style(el.get("style", ""))

    x = float(el.get("x", 0)) + tx
    y = float(el.get("y", 0)) + ty  # baseline position in SVG

    # Collect all text runs from direct text + tspans
    runs: list[dict] = []

    def _collect_runs(node: ET.Element, inherit: dict):
        nsty = {**inherit}
        child_sty = _parse_style(node.get("style", ""))
        nsty.update(child_sty)
        # Inline attrs override style
        for attr_name in ("fill", "font-size", "font-family", "font-weight",
                          "font-style", "text-anchor"):
            val = node.get(attr_name, "")
            if val:
                nsty[attr_name] = val

        # Direct text content (before any child tspans)
        if node.text and node.text.strip():
            runs.append({"text": node.text.strip(), "style": dict(nsty)})

        for child in node:
            tag = child.tag.split("}")[-1]
            if tag == "tspan":
                _collect_runs(child, nsty)
                if child.tail and child.tail.strip():
                    runs.append({"text": child.tail.strip(), "style": dict(nsty)})

    _collect_runs(el, {
        "fill": el.get("fill", sty.get("fill", "000000")),
        "font-size": el.get("font-size", sty.get("font-size", "16")),
        "font-family": el.get("font-family", sty.get("font-family", G.FONT_NAME)),
        "font-weight": el.get("font-weight", sty.get("font-weight", "normal")),
        "font-style": el.get("font-style", sty.get("font-style", "normal")),
        "text-anchor": el.get("text-anchor", sty.get("text-anchor", "start")),
    })

    if not runs:
        # No tspan — use element text directly
        txt = (el.text or "").strip()
        if not txt:
            return
        runs = [{"text": txt, "style": {
            "fill": el.get("fill", sty.get("fill", "000000")),
            "font-size": el.get("font-size", sty.get("font-size", "16")),
            "font-family": el.get("font-family", sty.get("font-family", G.FONT_NAME)),
            "font-weight": el.get("font-weight", sty.get("font-weight", "normal")),
            "font-style": el.get("font-style", sty.get("font-style", "normal")),
            "text-anchor": el.get("text-anchor", sty.get("text-anchor", "start")),
        }}]

    # Use first run's style for box sizing
    first_sty = runs[0]["style"]
    font_size_svg = _parse_font_size(first_sty.get("font-size", "16"))
    font_size_pt  = ctx.pt_font(font_size_svg)
    anchor        = first_sty.get("text-anchor", "start")

    # Estimate text box dimensions
    all_text = " ".join(r["text"] for r in runs)
    est_w = min(len(all_text) * font_size_pt * 0.58, _SLIDE_W_PT - x)
    est_h = font_size_pt * 2.0   # height: 2× font size for safety

    # Adjust x for text-anchor
    if anchor == "middle":
        x -= est_w / 2
    elif anchor == "end":
        x -= est_w

    # SVG y is baseline → top of textbox
    box_y = y - font_size_svg

    align_map = {"start": PP_ALIGN.LEFT, "middle": PP_ALIGN.CENTER, "end": PP_ALIGN.RIGHT}
    align = align_map.get(anchor, PP_ALIGN.LEFT)

    # Clamp to slide
    x = max(0.0, x)
    box_y = max(0.0, box_y)

    txBox = slide.shapes.add_textbox(
        Emu(ctx.emu_x(x)), Emu(ctx.emu_y(box_y)),
        Emu(ctx.emu_w(max(est_w, 40))), Emu(ctx.emu_h(max(est_h, 20))),
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    bodyPr.set("anchor", "t")
    bodyPr.set("lIns", "0")
    bodyPr.set("rIns", "0")
    bodyPr.set("tIns", "0")
    bodyPr.set("bIns", "0")
    for tag in (qn("a:normAutofit"), qn("a:noAutofit"), qn("a:spAutoFit")):
        for child in bodyPr.findall(tag):
            bodyPr.remove(child)
    lxml_etree.SubElement(bodyPr, qn("a:normAutofit"))

    p = tf.paragraphs[0]
    p.alignment = align

    for i, run_data in enumerate(runs):
        rsty = run_data["style"]
        fs_svg = _parse_font_size(rsty.get("font-size", "16"))
        fs_pt  = ctx.pt_font(fs_svg)
        color  = _parse_color(rsty.get("fill", "000000")) or "000000"
        bold   = rsty.get("font-weight", "normal").lower() in ("bold", "700", "800", "900")
        italic = rsty.get("font-style", "normal").lower() == "italic"
        font   = rsty.get("font-family", G.FONT_NAME).split(",")[0].strip().strip("'\"")

        run = p.add_run()
        run.text = run_data["text"] + (" " if i < len(runs) - 1 else "")
        run.font.name = font or G.FONT_NAME
        run.font.size = Pt(max(fs_pt, 6))
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = RGBColor(
            int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
        )


def _render_image(slide, el: ET.Element, ctx: _Ctx, tx: float, ty: float):
    """Render <image> or <use href="..."> elements."""
    x = float(el.get("x", 0)) + tx
    y = float(el.get("y", 0)) + ty
    w = float(el.get("width", 0))
    h = float(el.get("height", 0))
    if w <= 0 or h <= 0:
        return

    href = (
        el.get("href", "")
        or el.get(f"{{{_XLINK}}}href", "")
    )
    if not href:
        return

    try:
        if href.startswith("data:"):
            # data URI: data:image/png;base64,....
            m = re.match(r"data:[^;]+;base64,(.+)", href)
            if not m:
                return
            import base64
            img_bytes = base64.b64decode(m.group(1))
            img_stream = io.BytesIO(img_bytes)
        elif href.startswith(("http://", "https://")):
            with urllib.request.urlopen(href, timeout=8) as resp:  # nosec
                img_bytes = resp.read()
            img_stream = io.BytesIO(img_bytes)
        else:
            return  # local file paths skipped for security

        slide.shapes.add_picture(
            img_stream,
            Emu(ctx.emu_x(x)), Emu(ctx.emu_y(y)),
            Emu(ctx.emu_w(w)), Emu(ctx.emu_h(h)),
        )
    except Exception:
        # If image download fails, draw a placeholder rect
        shape = slide.shapes.add_shape(
            1,
            Emu(ctx.emu_x(x)), Emu(ctx.emu_y(y)),
            Emu(ctx.emu_w(w)), Emu(ctx.emu_h(h)),
        )
        spPr = shape._element.spPr
        _set_solid_fill_on_spPr(spPr, "DDDDDD")
        _remove_line(spPr)
        shape.text_frame.text = ""


# ─────────────────────────────────────────────────────────────────────────────
# Walk element tree
# ─────────────────────────────────────────────────────────────────────────────

def _walk(slide, el: ET.Element, ctx: _Ctx, tx: float, ty: float):
    """Recursively render all SVG elements."""
    tag = el.tag.split("}")[-1] if "}" in el.tag else el.tag

    if tag == "defs":
        return  # already parsed in pre-pass

    transform = el.get("transform", "")
    dtx, dty  = _parse_transform(transform)
    atx, aty  = tx + dtx, ty + dty

    if tag == "g":
        for child in el:
            _walk(slide, child, ctx, atx, aty)

    elif tag == "rect":
        _render_rect(slide, el, ctx, atx, aty)

    elif tag in ("circle", "ellipse"):
        _render_oval(slide, el, ctx, atx, aty)

    elif tag == "line":
        _render_line(slide, el, ctx, atx, aty)

    elif tag == "text":
        _render_text(slide, el, ctx, atx, aty)

    elif tag == "image":
        _render_image(slide, el, ctx, atx, aty)

    # polygon / polyline / path → silently skip (too complex for lossless mapping)


# ─────────────────────────────────────────────────────────────────────────────
# Gradient pre-pass
# ─────────────────────────────────────────────────────────────────────────────

def _parse_gradients(root: ET.Element) -> dict:
    """Pre-scan <defs><linearGradient> and return id→{from, to, angle} dict."""
    grads: dict = {}
    for defs in root.iter(f"{{{_SVG}}}defs"):
        for lg in defs.iter(f"{{{_SVG}}}linearGradient"):
            gid = lg.get("id", "")
            if not gid:
                continue

            x1 = _pct(lg.get("x1", "0"))
            y1 = _pct(lg.get("y1", "0"))
            x2 = _pct(lg.get("x2", "1"))
            y2 = _pct(lg.get("y2", "0"))
            angle = _svg_grad_angle(x1, y1, x2, y2)

            stops = []
            for stop in lg.iter(f"{{{_SVG}}}stop"):
                ssty = _parse_style(stop.get("style", ""))
                color_val = ssty.get("stop-color", stop.get("stop-color", "000000"))
                color = _parse_color(color_val) or "000000"
                stops.append(color)

            if len(stops) >= 2:
                grads[gid] = {"from": stops[0], "to": stops[-1], "angle": angle}

    return grads


# ─────────────────────────────────────────────────────────────────────────────
# Public API
# ─────────────────────────────────────────────────────────────────────────────

def parse_svg_to_slide(prs: Presentation, svg_string: str) -> None:
    """Parse svg_string and append one slide to prs.

    Parameters
    ----------
    prs : Presentation
        The python-pptx Presentation object (already sized).
    svg_string : str
        A complete SVG document string.
    """
    # Normalise namespace — allow SVG without explicit namespace
    svg_string = svg_string.strip()
    if "<svg" in svg_string and _SVG not in svg_string:
        svg_string = svg_string.replace("<svg", f'<svg xmlns="{_SVG}"', 1)

    root = ET.fromstring(svg_string)

    # Parse viewBox to compute scale
    vb = root.get("viewBox", "")
    if vb:
        parts = re.split(r"[\s,]+", vb.strip())
        if len(parts) == 4:
            vb_w = float(parts[2])
            vb_h = float(parts[3])
        else:
            vb_w, vb_h = _SLIDE_W_PT, _SLIDE_H_PT
    else:
        # Try width/height attributes
        try:
            vb_w = float(re.sub(r"[^\d.]", "", root.get("width", str(_SLIDE_W_PT))))
            vb_h = float(re.sub(r"[^\d.]", "", root.get("height", str(_SLIDE_H_PT))))
        except ValueError:
            vb_w, vb_h = _SLIDE_W_PT, _SLIDE_H_PT

    scale_x = _SLIDE_W_PT / vb_w if vb_w else 1.0
    scale_y = _SLIDE_H_PT / vb_h if vb_h else 1.0

    gradients = _parse_gradients(root)
    ctx = _Ctx(scale_x, scale_y, gradients)

    # Add blank slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Walk top-level children (skip <defs>)
    for child in root:
        _walk(slide, child, ctx, 0.0, 0.0)
