# engine.py — Render orchestrator
from __future__ import annotations

import io
from pptx import Presentation
from pptx.util import Emu

from . import geometry as G
from .validator import validate
from .design_system import resolve_deck
from .html_renderer import VISUAL_LAYOUTS, render_to_png
from .layouts import (
    cover, card_grid, two_col_list, flow_steps, two_col_contrast,
    data_highlight, narrative, section_divider, cta, toc, timeline,
    icon_grid, quote
)

LAYOUT_MAP = {
    "I": cover.render,
    "A": card_grid.render,
    "B": two_col_list.render,
    "C": flow_steps.render,
    "D": two_col_contrast.render,
    "E": data_highlight.render,
    "F": narrative.render,
    "G": section_divider.render,
    "H": cta.render,
    "J": toc.render,
    "K": timeline.render,
    "L": icon_grid.render,
    "M": quote.render,
}


class ValidationError(Exception):
    def __init__(self, errors: list[str]):
        self.errors = errors
        super().__init__("Validation failed")


def _is_intent_deck(deck_json: dict) -> bool:
    """True when slides have no "design" key — i.e. semantic intent JSON from Gemini."""
    slides = deck_json.get("slides", [])
    return bool(slides) and not any("design" in s for s in slides)


def render_deck(deck_json: dict) -> bytes:
    """
    1. Auto-resolve semantic intent JSON → full design JSON if needed
    2. Validate JSON against schema — raise ValidationError if invalid
    3. Create Presentation with SLIDE_W × SLIDE_H
    4. For each slide: call LAYOUT_MAP[layout](prs, slide_data)
    5. Each layout calls inject_footer as its final step
    6. Save to BytesIO, return bytes
    """
    if _is_intent_deck(deck_json):
        deck_json = resolve_deck(deck_json)

    ok, errors = validate(deck_json)
    if not ok:
        raise ValidationError(errors)

    prs = Presentation()
    prs.slide_width = Emu(G.SLIDE_W)
    prs.slide_height = Emu(G.SLIDE_H)

    for slide_data in deck_json["slides"]:
        layout_key = slide_data["layout"]

        if layout_key in VISUAL_LAYOUTS:
            # Render via Playwright → embed as full-slide PNG picture
            png_bytes = render_to_png(layout_key, slide_data)
            blank_layout = prs.slide_layouts[6]  # blank master
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                io.BytesIO(png_bytes),
                Emu(0), Emu(0),
                Emu(G.SLIDE_W), Emu(G.SLIDE_H),
            )
        else:
            renderer = LAYOUT_MAP.get(layout_key)
            if renderer is None:
                raise ValueError(f"Unknown layout: {layout_key!r}")
            renderer(prs, slide_data)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
