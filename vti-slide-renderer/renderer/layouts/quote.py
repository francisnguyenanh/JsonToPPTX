# layouts/quote.py — Layout M (Quote / Key Message)
from __future__ import annotations

from pptx.util import Emu
from lxml import etree
from pptx.oxml.ns import qn

from .. import geometry as G
from ..shapes import (
    add_textbox_styled, add_rounded_rect,
    _set_gradient_fill_on_spPr, _set_solid_fill_on_spPr,
)
from ..footer import inject_footer
from ._common import add_slide_bg, add_accent_bar, add_slide_header, remove_line


def render(prs, slide_data: dict):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    d = slide_data["design"]
    sn = slide_data["slide_number"]

    add_slide_bg(slide, d.get("bg", {"from": "FFFFFF", "to": "EEF4FB"}))
    add_accent_bar(slide, d.get("accent_bar", {"from": "2362B0", "to": "7FC236"}))
    add_slide_header(slide, d.get("title", ""), d.get("breadcrumb", ""),
                     title_color=d.get("title_color", "172759"),
                     breadcrumb_color=d.get("breadcrumb_color", "6A7FA0"))

    # ── Left accent block (16pt wide gradient bar) ───────────────────
    alb = d.get("accent_block_left", {"from": "2362B0", "to": "7FC236"})
    bar_w = G.pt(19)
    accent_block = slide.shapes.add_shape(
        1, Emu(G.CONTENT_X), Emu(G.CONTENT_TOP),
        Emu(bar_w), Emu(G.CONTENT_H - G.pt(87))
    )
    _set_gradient_fill_on_spPr(accent_block._element.spPr,
                                alb["from"], alb["to"], alb.get("angle", 0))
    remove_line(accent_block._element.spPr)
    accent_block.text_frame.text = ""

    # ── Decorative large quote mark ────────────────────────────────────
    qmark_color = d.get("quote_mark_color", "2362B0")
    add_textbox_styled(
        slide,
        G.CONTENT_X + bar_w, G.CONTENT_TOP - G.pt(13),
        G.pt(107), G.pt(96), "“",
        size_pt=120, color_hex=qmark_color,
        bold=True, v_anchor="t", inset=G.INS_NONE, autofit="none"
    )

    # ── Quote text zone — starts below the quote mark glyph ──────────
    quote_x = G.CONTENT_X + bar_w + G.pt(48)
    quote_w = G.CONTENT_W - bar_w - G.pt(61)
    quote_top = G.CONTENT_TOP + G.pt(96)
    context_h = G.pt(77)
    quote_h = G.CONTENT_H - G.pt(96) - context_h - G.pt(11)

    # quote_text is the primary key; fall back to legacy "quote"
    quote_text = d.get("quote_text", d.get("quote", ""))
    add_textbox_styled(
        slide,
        quote_x, quote_top, quote_w, quote_h,
        quote_text,
        size_pt=20, bold=False, italic=True,
        color_hex=d.get("quote_color", "172759"),
        align="left", v_anchor="m",
        inset=G.INS_CARD, autofit="norm", wrap=True,
        line_spacing=1.5
    )

    # Attribution
    if d.get("attribution"):
        attr_y = G.CONTENT_TOP + G.CONTENT_H - context_h - G.pt(40)
        attr_text = d["attribution"]
        if not attr_text.startswith("—"):
            attr_text = "— " + attr_text
        add_textbox_styled(
            slide,
            quote_x, attr_y, quote_w, G.pt(32),
            attr_text,
            size_pt=G.FONT_BODY,
            color_hex=d.get("attribution_color", "6A7FA0"),
            align="right", v_anchor="m",
            inset=G.INS_NONE, autofit="none"
        )

    # ── Context card (bottom strip) ──────────────────────────────────
    ctx_text = d.get("context_text", "")
    if ctx_text:
        ctx_y = G.CONTENT_TOP + G.CONTENT_H - context_h
        ctx_bg = d.get("context_bg", {"from": "EEF4FB", "to": "FFFFFF"})
        add_rounded_rect(slide, G.CONTENT_X, ctx_y, G.CONTENT_W, context_h,
                         ctx_bg["from"], ctx_bg["to"],
                         angle_emu=ctx_bg.get("angle", 16200000))
        add_textbox_styled(
            slide,
            G.CONTENT_X + G.pt(16), ctx_y,
            G.CONTENT_W - G.pt(32), context_h,
            ctx_text,
            size_pt=G.FONT_BODY,
            color_hex=d.get("context_text_color", "1C2D4F"),
            v_anchor="m", inset=G.INS_NONE, autofit="none", wrap=True
        )

    inject_footer(slide, sn)
