# layouts/card_grid.py — Layout A
from __future__ import annotations

from pptx.util import Emu
from .. import geometry as G
from ..shapes import (
    add_textbox_styled, add_rounded_rect, add_dot_bullet,
    add_semantic_icon, add_separator_line, _set_solid_fill_on_spPr
)
from ..footer import inject_footer
from ._common import add_slide_bg, add_accent_bar, add_slide_header, remove_line
from lxml import etree
from pptx.oxml.ns import qn


def render(prs, slide_data: dict):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    d = slide_data["design"]
    sn = slide_data["slide_number"]

    add_slide_bg(slide, d.get("bg", {"from": "FFFFFF", "to": "EEF4FB"}))
    add_accent_bar(slide, d.get("accent_bar", {"from": "2362B0", "to": "7FC236"}))
    add_slide_header(slide, d.get("title", ""), d.get("breadcrumb", ""))

    cards = d.get("cards", [])
    n_cols = d.get("n_cols", len(cards))
    n_rows = 1

    # Adaptive typography based on column count
    if n_cols <= 2:
        header_font = 20
        body_font = 15
        header_h = G.pt(52)
        header_adv = G.pt(56)
        bullet_h = G.pt(24)
        bullet_adv = G.pt(26)
        icon_sz_default = 38
        pad = G.pt(16)
    elif n_cols == 3:
        header_font = 17
        body_font = 14
        header_h = G.pt(46)
        header_adv = G.pt(50)
        bullet_h = G.pt(34)
        bullet_adv = G.pt(36)
        icon_sz_default = 32
        pad = G.pt(14)
    else:  # 4 columns
        header_font = 15
        body_font = 13
        header_h = G.pt(42)
        header_adv = G.pt(46)
        bullet_h = G.pt(30)
        bullet_adv = G.pt(32)
        icon_sz_default = 28
        pad = G.pt(12)

    for col_idx, card in enumerate(cards):
        cx = G.card_x(col_idx, n_cols)
        cy = G.card_y(0, n_rows)
        cw = G.card_w(n_cols)
        ch = G.card_h(n_rows)

        card_bg = card.get("bg", {"from": "EBF4FC", "to": "D6EAF8"})
        add_rounded_rect(slide, cx, cy, cw, ch,
                         card_bg["from"], card_bg["to"],
                         angle_emu=16200000,
                         border_hex=card.get("border_top"),
                         border_w_pt=3)

        inner_x = cx + pad
        inner_w = cw - 2 * pad
        cur_y = cy + pad

        # Icon
        icon = card.get("icon")
        if icon:
            sz = icon.get("size_pt", icon_sz_default)
            add_semantic_icon(slide, inner_x, cur_y, sz,
                              icon.get("stroke", "4A9EE0"),
                              icon.get("type", "default"))
            cur_y += G.pt(sz + 10)

        # Header
        if card.get("header"):
            add_textbox_styled(slide, inner_x, cur_y, inner_w, header_h,
                               card["header"], bold=True,
                               size_pt=header_font,
                               color_hex=card.get("header_color", "172759"),
                               v_anchor="t", inset=G.INS_NONE,
                               autofit="norm", wrap=True)
            cur_y += header_adv

        # Thin separator line below header
        if card.get("header") and card.get("bullets"):
            sep_color = card.get("border_top", "4A9EE0")
            add_separator_line(slide, inner_x, cur_y,
                               inner_w * 2 // 5, G.pt(1), sep_color)
            cur_y += G.pt(8)

        # Bullets
        for bullet in card.get("bullets", []):
            dot_color = bullet.get("dot_color", "4A9EE0")
            dot_y = cur_y + G.pt(4)
            add_dot_bullet(slide, inner_x, dot_y, dot_color, 5)
            add_textbox_styled(slide, inner_x + G.pt(12), cur_y,
                               inner_w - G.pt(12), bullet_h,
                               bullet.get("text", ""),
                               size_pt=body_font,
                               color_hex=card.get("text_color", "1C2D4F"),
                               v_anchor="t", inset=G.INS_NONE,
                               autofit="norm", wrap=True)
            cur_y += bullet_adv

    inject_footer(slide, sn)
