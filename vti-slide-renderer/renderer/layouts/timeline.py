# layouts/timeline.py  ELayout K
from __future__ import annotations

from pptx.util import Emu
from .. import geometry as G
from ..shapes import (
    add_textbox_styled, add_rounded_rect, add_dot_bullet,
    _set_gradient_fill_on_spPr
)
from ..footer import inject_footer
from ._common import add_slide_bg, add_accent_bar, add_slide_header, remove_line


def render(prs, slide_data: dict):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    d = slide_data["design"]
    sn = slide_data["slide_number"]

    add_slide_bg(slide, d.get("bg", {"from": "FFFFFF", "to": "EEF4FB"}))
    add_accent_bar(slide, d.get("accent_bar", {"from": "2362B0", "to": "7FC236"}))
    add_slide_header(slide, d.get("title", ""), d.get("breadcrumb", ""))

    events = d.get("events", [])
    n = len(events)
    spine = d.get("spine", {"from": "2362B0", "to": "7FC236"})

    # Spine (horizontal gradient bar in middle of content zone)
    spine_y = G.CONTENT_TOP + G.CONTENT_H // 2 - G.pt(1)
    spine_shape = slide.shapes.add_shape(
        1, Emu(G.CONTENT_X), Emu(spine_y), Emu(G.CONTENT_W), Emu(G.pt(2))
    )
    _set_gradient_fill_on_spPr(spine_shape._element.spPr,
                                spine["from"], spine["to"], 5400000)
    remove_line(spine_shape._element.spPr)
    spine_shape.text_frame.text = ""

    event_w = G.card_w(n)
    card_h = G.CONTENT_H // 2 - G.pt(21)

    for i, event in enumerate(events):
        ex = G.card_x(i, n)
        bg = event.get("bg", {"from": "EBF4FC", "to": "D6EAF8"})

        # Alternate above/below spine
        if i % 2 == 0:
            ey = G.CONTENT_TOP
        else:
            ey = spine_y + G.pt(11)

        add_rounded_rect(slide, ex, ey, event_w, card_h,
                         bg["from"], bg["to"], shadow=True)

        # Dot on spine
        dot_x = ex + event_w // 2 - G.pt(7)
        dot_y = spine_y - G.pt(5)
        add_dot_bullet(slide, dot_x, dot_y, event.get("dot_color", "2362B0"), 10)

        # Date
        add_textbox_styled(slide, ex + G.pt(11), ey + G.pt(8),
                           event_w - G.pt(21), G.pt(24),
                           event.get("date", ""), bold=True, size_pt=G.FONT_SMALL,
                           color_hex=event.get("date_color", "6A7FA0"),
                           v_anchor="t", inset=G.INS_NONE, autofit="none")

        # Label
        add_textbox_styled(slide, ex + G.pt(11), ey + G.pt(32),
                           event_w - G.pt(21), G.pt(29),
                           event.get("label", ""), bold=True, size_pt=G.FONT_BODY,
                           color_hex=event.get("text_color", "1C2D4F"),
                           v_anchor="t", inset=G.INS_NONE, autofit="none", wrap=True)

        # Description
        if event.get("description"):
            add_textbox_styled(slide, ex + G.pt(11), ey + G.pt(67),
                               event_w - G.pt(21), card_h - G.pt(75),
                               event["description"], size_pt=G.FONT_SMALL,
                               color_hex=event.get("text_color", "1C2D4F"),
                               v_anchor="t", inset=G.INS_NONE, autofit="norm", wrap=True)

    inject_footer(slide, sn)
