# layouts/flow_steps.py — Layout C
from __future__ import annotations

from pptx.util import Emu
from .. import geometry as G
from ..shapes import add_textbox_styled, add_rounded_rect, add_badge, add_separator_line
from ..footer import inject_footer
from ._common import add_slide_bg, add_accent_bar, add_slide_header


def render(prs, slide_data: dict):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    d = slide_data["design"]
    sn = slide_data["slide_number"]

    add_slide_bg(slide, d.get("bg", {"from": "FFFFFF", "to": "EEF4FB"}))
    add_accent_bar(slide, d.get("accent_bar", {"from": "2362B0", "to": "7FC236"}))
    add_slide_header(slide, d.get("title", ""), d.get("breadcrumb", ""))

    steps = d.get("steps", [])
    n = len(steps)
    direction = d.get("flow_direction", "horizontal")
    arrow_color = d.get("arrow_color", "B8D0EE")

    if direction == "horizontal":
        step_w = G.card_w(n)
        step_h = G.CONTENT_H

        # Arrow connector dimensions
        arr_gap = G.COL_GAP

        for i, step in enumerate(steps):
            sx = G.card_x(i, n)
            sy = G.CONTENT_TOP
            bg = step.get("bg", {"from": "F4F9FE", "to": "EAF3FB"})
            add_rounded_rect(slide, sx, sy, step_w, step_h,
                             bg["from"], bg["to"],
                             shadow=True)

            # Arrow between steps — centered vertically
            if i < n - 1:
                arr_x = sx + step_w + G.pt(2)
                arr_y = sy + step_h // 2 - G.pt(1)
                add_separator_line(slide, arr_x, arr_y,
                                   arr_gap - G.pt(4), G.pt(2), arrow_color)

            inner_x = sx + G.pt(19)
            inner_w = step_w - G.pt(37)

            # Badge size
            badge_sz = 21 if n >= 4 else 24
            badge_block = G.pt(badge_sz + 8)

            # Header height — wider for fewer columns
            hdr_h = G.pt(69) if n <= 3 else G.pt(64)
            hdr_adv = G.pt(75) if n <= 3 else G.pt(69)

            # Estimate content height for vertical centering
            has_desc = bool(step.get("description"))
            desc_h_est = G.pt(48) if has_desc else 0
            total_est = badge_block + hdr_adv + desc_h_est
            avail = step_h - G.pt(37)
            offset = max(G.pt(21), (avail - total_est) // 2)
            cur_y = sy + offset

            add_badge(slide, inner_x, cur_y, str(i + 1),
                      step.get("badge_color", "2362B0"), size_pt=badge_sz)
            cur_y += badge_block

            if step.get("header"):
                add_textbox_styled(slide, inner_x, cur_y, inner_w, hdr_h,
                                   step["header"], bold=True,
                                   size_pt=G.FONT_HEADER if n <= 3 else 17,
                                   color_hex=step.get("header_color", "172759"),
                                   v_anchor="t", inset=G.INS_NONE,
                                   autofit="norm", wrap=True)
                cur_y += hdr_adv

            if step.get("description"):
                desc_h = step_h - (cur_y - sy) - G.pt(19)
                add_textbox_styled(slide, inner_x, cur_y, inner_w, desc_h,
                                   step["description"],
                                   size_pt=G.FONT_BODY if n <= 3 else 15,
                                   color_hex=step.get("text_color", "1C2D4F"),
                                   v_anchor="t", inset=G.INS_NONE,
                                   autofit="norm", wrap=True)
    else:
        # Vertical flow
        step_h = G.card_h(n)
        for i, step in enumerate(steps):
            sx = G.CONTENT_X
            sy = G.card_y(i, n)
            bg = step.get("bg", {"from": "F4F9FE", "to": "EAF3FB"})
            add_rounded_rect(slide, sx, sy, G.CONTENT_W, step_h,
                             bg["from"], bg["to"],
                             shadow=True)

            inner_x = sx + G.pt(19)
            cur_y = sy + G.pt(16)
            add_badge(slide, inner_x, cur_y, str(i + 1),
                      step.get("badge_color", "2362B0"), size_pt=21)

            if step.get("header"):
                add_textbox_styled(slide, inner_x + G.pt(29), cur_y,
                                   G.CONTENT_W - G.pt(67), G.pt(40),
                                   step["header"], bold=True,
                                   size_pt=G.FONT_HEADER,
                                   color_hex=step.get("header_color", "172759"),
                                   v_anchor="m", inset=G.INS_NONE,
                                   autofit="norm", wrap=True)

    inject_footer(slide, sn)
