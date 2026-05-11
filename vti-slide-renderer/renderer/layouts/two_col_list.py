# layouts/two_col_list.py — Layout B
from __future__ import annotations

from pptx.util import Emu
from .. import geometry as G
from ..shapes import add_textbox_styled, add_dot_bullet, add_separator_line
from ..footer import inject_footer
from ._common import add_slide_bg, add_accent_bar, add_slide_header


def render(prs, slide_data: dict):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    d = slide_data["design"]
    sn = slide_data["slide_number"]

    add_slide_bg(slide, d.get("bg", {"from": "FFFFFF", "to": "EEF4FB"}))
    add_accent_bar(slide, d.get("accent_bar", {"from": "2362B0", "to": "7FC236"}))
    add_slide_header(slide, d.get("title", ""), d.get("breadcrumb", ""))

    col_w = (G.CONTENT_W - G.pt(1)) // 2
    left_x = G.CONTENT_X
    right_x = G.CONTENT_X + col_w + G.pt(1)

    # Vertical separator
    sep_color = d.get("separator_color", "B8D0EE")
    add_separator_line(slide, G.CONTENT_X + col_w, G.CONTENT_TOP,
                       G.pt(1), G.CONTENT_H, sep_color)

    def draw_items(items, start_x):
        cy = G.CONTENT_TOP + G.pt(11)
        row_h = G.pt(69)
        for item in items:
            dot_color = item.get("badge_color", "2362B0")
            dot_size = 10
            # Vertically centre the dot with the text line
            add_dot_bullet(slide, start_x + G.pt(5),
                           cy + G.pt(row_h // 2 // G.PT - dot_size // 2),
                           dot_color, dot_size)
            add_textbox_styled(slide, start_x + G.pt(24), cy,
                               col_w - G.pt(27), row_h,
                               item.get("text", ""),
                               size_pt=G.FONT_BODY,
                               color_hex=item.get("text_color", "1C2D4F"),
                               v_anchor="m", inset=G.INS_NONE, autofit="norm", wrap=True)
            cy += row_h + G.pt(5)

    draw_items(d.get("left_items", []), left_x)
    draw_items(d.get("right_items", []), right_x)

    inject_footer(slide, sn)
