# layouts/data_highlight.py — Layout E
from __future__ import annotations

from pptx.util import Emu
from .. import geometry as G
from ..shapes import add_textbox_styled, add_rounded_rect, add_separator_line
from ..footer import inject_footer
from ._common import add_slide_bg, add_accent_bar, add_slide_header


def render(prs, slide_data: dict):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    d = slide_data["design"]
    sn = slide_data["slide_number"]

    add_slide_bg(slide, d.get("bg", {"from": "FFFFFF", "to": "EEF4FB"}))
    add_accent_bar(slide, d.get("accent_bar", {"from": "2362B0", "to": "7FC236"}))
    add_slide_header(slide, d.get("title", ""), d.get("breadcrumb", ""))

    stats = d.get("stats", [])
    n = len(stats)

    # Stats area: top ~60% of content zone
    stats_h = G.CONTENT_H * 58 // 100
    stat_w = G.card_w(n)

    for i, stat in enumerate(stats):
        sx = G.card_x(i, n)
        sy = G.CONTENT_TOP
        bg = stat.get("bg", {"from": "F4F9FE", "to": "EAF3FB"})
        add_rounded_rect(slide, sx, sy, stat_w, stats_h,
                         bg["from"], bg["to"], shadow=True)

        inner_x = sx + G.pt(13)
        inner_w = stat_w - G.pt(27)

        # Number — prominent, top-aligned, normAutofit shrinks if too wide
        num_str = stat.get("number", "")
        add_textbox_styled(slide, inner_x, sy + G.pt(21), inner_w, G.pt(113),
                           num_str,
                           bold=True, size_pt=G.FONT_STAT_BIG,
                           color_hex=stat.get("number_color", "2362B0"),
                           align="center", v_anchor="m",
                           inset=G.INS_NONE, autofit="norm")

        # Thin separator line
        sep_color = stat.get("number_color", "2362B0")
        sep_w = inner_w * 2 // 5
        sep_x = sx + (stat_w - sep_w) // 2
        add_separator_line(slide, sep_x, sy + G.pt(141), sep_w, G.pt(3), sep_color)

        # Label — below separator, 2-line max, normAutofit
        add_textbox_styled(slide, inner_x, sy + G.pt(149), inner_w, G.pt(77),
                           stat.get("label", ""),
                           size_pt=G.FONT_STAT_LABEL,
                           color_hex=stat.get("label_color", "6A7FA0"),
                           align="center", v_anchor="t",
                           inset=G.INS_NONE, autofit="norm", wrap=True)

    # Context strip
    ctx = d.get("context_text", "")
    if ctx:
        ctx_y = G.CONTENT_TOP + stats_h + G.CARD_GAP
        ctx_h = G.CONTENT_H - stats_h - G.CARD_GAP
        ctx_bg = d.get("context_bg", {"from": "EEF4FB", "to": "FFFFFF"})
        add_rounded_rect(slide, G.CONTENT_X, ctx_y, G.CONTENT_W, ctx_h,
                         ctx_bg["from"], ctx_bg["to"])
        add_textbox_styled(slide, G.CONTENT_X + G.pt(21), ctx_y,
                           G.CONTENT_W - G.pt(43), ctx_h,
                           ctx,
                           size_pt=G.FONT_BODY,
                           color_hex=d.get("context_text_color", "1C2D4F"),
                           v_anchor="m", inset=G.INS_NONE, autofit="norm", wrap=True)

    inject_footer(slide, sn)
