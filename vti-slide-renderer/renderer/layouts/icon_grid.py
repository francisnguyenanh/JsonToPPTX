# layouts/icon_grid.py — Layout L
from __future__ import annotations

from pptx.util import Emu
from .. import geometry as G
from ..shapes import add_textbox_styled, add_rounded_rect, add_semantic_icon
from ..footer import inject_footer
from ._common import add_slide_bg, add_accent_bar, add_slide_header


def render(prs, slide_data: dict):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    d = slide_data["design"]
    sn = slide_data["slide_number"]

    add_slide_bg(slide, d.get("bg", {"from": "FFFFFF", "to": "EEF4FB"}))
    add_accent_bar(slide, d.get("accent_bar", {"from": "2362B0", "to": "7FC236"}))
    add_slide_header(slide, d.get("title", ""), d.get("breadcrumb", ""))

    cells = d.get("cells", [])
    n_cols = d.get("n_cols", 3)
    n_rows = (len(cells) + n_cols - 1) // n_cols

    cw = G.card_w(n_cols)
    ch = G.card_h(n_rows)

    for idx, cell in enumerate(cells):
        col = idx % n_cols
        row = idx // n_cols
        cx = G.card_x(col, n_cols)
        cy = G.card_y(row, n_rows)

        bg = cell.get("bg", {"from": "EBF4FC", "to": "D6EAF8"})
        add_rounded_rect(slide, cx, cy, cw, ch,
                         bg["from"], bg["to"],
                         border_hex=cell.get("border_top"), border_w_pt=3)

        inner_x = cx + G.pt(8)
        inner_w = cw - G.pt(16)

        # Calculate vertical layout: icon + label + sub_label
        icon = cell.get("icon")
        has_label = bool(cell.get("label"))
        has_sub = bool(cell.get("sub_label"))
        icon_sz = icon.get("size_pt", 36) if icon else 0
        label_h = G.pt(24) if has_label else 0
        sub_h = G.pt(20) if has_sub else 0
        gap = G.pt(6)
        total_h = G.pt(icon_sz) + (gap if icon_sz else 0) + label_h + sub_h
        cur_y = cy + (ch - total_h) // 2  # vertical center

        if icon:
            # Horizontally centre icon
            icon_x = cx + (cw - G.pt(icon_sz)) // 2
            add_semantic_icon(slide, icon_x, cur_y, icon_sz,
                              icon.get("stroke", "4A9EE0"),
                              icon.get("type", "default"))
            cur_y += G.pt(icon_sz) + gap

        if has_label:
            add_textbox_styled(slide, inner_x, cur_y, inner_w, label_h,
                               cell["label"], bold=True,
                               size_pt=G.FONT_CARD_HEADER,
                               color_hex=cell.get("label_color", "172759"),
                               align="center", v_anchor="m",
                               inset=G.INS_NONE, autofit="norm", wrap=True)
            cur_y += label_h

        if has_sub:
            add_textbox_styled(slide, inner_x, cur_y, inner_w, sub_h,
                               cell["sub_label"], size_pt=G.FONT_SMALL,
                               color_hex=cell.get("sub_label_color", "6A7FA0"),
                               align="center", v_anchor="t",
                               inset=G.INS_NONE, autofit="norm", wrap=True)

    inject_footer(slide, sn)
