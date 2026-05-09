# layouts/_common.py  EShared helpers for content layouts
from __future__ import annotations

from pptx.util import Emu
from lxml import etree
from pptx.oxml.ns import qn

from .. import geometry as G
from ..shapes import (
    _set_gradient_fill_on_spPr, _set_solid_fill_on_spPr,
    add_textbox_styled, add_separator_line
)


def remove_line(spPr):
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        spPr.remove(ln)
    etree.SubElement(spPr, qn("a:ln")).append(
        etree.fromstring(
            '<a:noFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
        )
    )


def add_slide_bg(slide, bg: dict):
    """Full-slide gradient background."""
    shape = slide.shapes.add_shape(
        1, Emu(0), Emu(0), Emu(G.SLIDE_W), Emu(G.SLIDE_H)
    )
    spPr = shape._element.spPr
    _set_gradient_fill_on_spPr(spPr, bg["from"], bg["to"], bg.get("angle", 9000000))
    remove_line(spPr)
    shape.text_frame.text = ""
    return shape


def add_accent_bar(slide, bar: dict):
    """Horizontal accent bar just below the title area."""
    shape = slide.shapes.add_shape(
        1, Emu(G.CONTENT_X), Emu(G.ACCENT_BAR_Y),
        Emu(G.CONTENT_W), Emu(G.ACCENT_BAR_H)
    )
    spPr = shape._element.spPr
    _set_gradient_fill_on_spPr(spPr, bar["from"], bar["to"], bar.get("angle", 5400000))
    remove_line(spPr)
    shape.text_frame.text = ""
    return shape


def add_slide_header(slide, title: str, breadcrumb: str = ""):
    """Slide title + breadcrumb."""
    add_textbox_styled(
        slide,
        G.CONTENT_X, G.pt(10),
        G.CONTENT_W, G.pt(50),
        title,
        bold=True, size_pt=G.FONT_TITLE,
        color_hex="172759",
        v_anchor="m", inset=G.INS_NONE, autofit="none"
    )
    if breadcrumb:
        add_textbox_styled(
            slide,
            G.CONTENT_X, G.pt(60),
            G.CONTENT_W, G.pt(24),
            breadcrumb,
            size_pt=G.FONT_SMALL, color_hex="6A7FA0",
            v_anchor="m", inset=G.INS_NONE, autofit="none"
        )
