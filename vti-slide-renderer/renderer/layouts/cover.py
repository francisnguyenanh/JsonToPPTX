# layouts/cover.py — Layout I
from __future__ import annotations

from pptx.util import Emu
from lxml import etree
from pptx.oxml.ns import qn

from .. import geometry as G
from ..shapes import (
    _set_gradient_fill_on_spPr, _set_solid_fill_on_spPr,
    add_textbox_styled, add_semantic_icon, add_rounded_rect,
    add_decor_shape,
)
from ..footer import inject_footer


def render(prs, slide_data: dict):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    d = slide_data["design"]
    sn = slide_data["slide_number"]

    # ── Background gradient ──────────────────────────────────────────
    bg = d.get("bg", {"from": "FFFFFF", "to": "EEF4FB"})
    bg_shape = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(G.SLIDE_W), Emu(G.SLIDE_H))
    _set_gradient_fill_on_spPr(bg_shape._element.spPr,
                                bg["from"], bg["to"], bg.get("angle", 9000000))
    _rm(bg_shape._element.spPr)
    bg_shape.text_frame.text = ""

    # ── Left vertical bar ────────────────────────────────────────────
    lb = d.get("left_bar", {"from": "172759", "to": "2362B0"})
    bar_w = G.pt(8)
    left_bar = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(bar_w), Emu(G.SLIDE_H))
    _set_gradient_fill_on_spPr(left_bar._element.spPr,
                                lb["from"], lb["to"], lb.get("angle", 0))
    _rm(left_bar._element.spPr)
    left_bar.text_frame.text = ""

    # ── Right decorative strip (semi-transparent) ────────────────────
    rs = d.get("right_strip", {"from": "2362B0", "to": "4A9EE0", "opacity": 30})
    strip_w = G.pt(140)
    strip_x = G.SLIDE_W - strip_w
    opacity = rs.get("opacity", 30)
    alpha_emu = int(opacity / 100 * 100000)
    rstrip = slide.shapes.add_shape(1, Emu(strip_x), Emu(0), Emu(strip_w), Emu(G.SLIDE_H))
    rs_spPr = rstrip._element.spPr
    for tag in (qn("a:noFill"), qn("a:solidFill"), qn("a:gradFill")):
        for el in rs_spPr.findall(tag):
            rs_spPr.remove(el)
    gf = etree.SubElement(rs_spPr, qn("a:gradFill"))
    gsl = etree.SubElement(gf, qn("a:gsLst"))
    for pos, hx in [(0, rs["from"]), (100000, rs["to"])]:
        gs = etree.SubElement(gsl, qn("a:gs"), pos=str(pos))
        sc = etree.SubElement(gs, qn("a:srgbClr"), val=hx.upper())
        etree.SubElement(sc, qn("a:alpha"), val=str(alpha_emu))
    etree.SubElement(gf, qn("a:lin"), ang=str(rs.get("angle", 0)), scaled="0")
    _rm(rs_spPr)
    rstrip.text_frame.text = ""

    # ── Right card (tinted, roundRect with top border + hero icon) ───
    rc = d.get("right_card", {})
    card_x = strip_x - G.pt(260)
    card_y = G.pt(60)
    card_w = G.pt(230)
    card_h = G.pt(410)

    rc_bg = rc.get("bg", {"from": "EBF4FC", "to": "D6EAF8"})
    add_rounded_rect(slide, card_x, card_y, card_w, card_h,
                     rc_bg["from"], rc_bg["to"],
                     angle_emu=rc_bg.get("angle", 16200000),
                     border_hex=rc.get("border_top"),
                     border_w_pt=4)

    # Decor ellipse inside right card (Pattern 8 — subtle depth)
    add_decor_shape(slide,
                    card_x + G.pt(100), card_y + G.pt(260),
                    G.pt(180), G.pt(180),
                    rc.get("border_top", "7FC236"), alpha_percent=7)

    # Hero icon centred in right card
    icon_def = rc.get("icon", {})
    if icon_def:
        icon_sz = icon_def.get("size_pt", 80)
        icon_x = card_x + (card_w - G.pt(icon_sz)) // 2
        icon_y = card_y + (card_h - G.pt(icon_sz)) // 2
        add_semantic_icon(slide, icon_x, icon_y, icon_sz,
                          icon_def.get("stroke", "4A9EE0"),
                          icon_def.get("type", "default"))

    # ── Divider bar ──────────────────────────────────────────────────
    div = d.get("divider_bar", {"from": "2362B0", "to": "7FC236"})
    tx = bar_w + G.pt(20)
    text_w = card_x - tx - G.pt(20)

    div_y = G.pt(258)
    div_shape = slide.shapes.add_shape(1, Emu(tx), Emu(div_y), Emu(G.pt(200)), Emu(G.pt(4)))
    _set_gradient_fill_on_spPr(div_shape._element.spPr,
                                div["from"], div["to"], 5400000)
    _rm(div_shape._element.spPr)
    div_shape.text_frame.text = ""

    # ── Text content (left zone) ─────────────────────────────────────
    # Tag
    if d.get("tag"):
        add_textbox_styled(slide, tx, G.pt(78), text_w, G.pt(26),
                           d["tag"], bold=True,
                           size_pt=12, color_hex=d.get("tag_color", "2362B0"),
                           v_anchor="m", inset=G.INS_NONE, autofit="none")

    # Title
    add_textbox_styled(slide, tx, G.pt(106), text_w, G.pt(145),
                       d.get("title", ""), bold=True,
                       size_pt=G.FONT_COVER_TITLE,
                       color_hex=d.get("title_color", "172759"),
                       v_anchor="t", inset=G.INS_NONE, autofit="norm", wrap=True)

    # Subtitle
    if d.get("subtitle"):
        add_textbox_styled(slide, tx, G.pt(272), text_w, G.pt(55),
                           d["subtitle"], size_pt=G.FONT_COVER_SUB,
                           color_hex=d.get("subtitle_color", "172759"),
                           v_anchor="m", inset=G.INS_NONE, autofit="none", wrap=True)

    # Caption
    if d.get("caption"):
        add_textbox_styled(slide, tx, G.pt(340), text_w, G.pt(28),
                           d["caption"], size_pt=12,
                           color_hex=d.get("caption_color", "6A7FA0"),
                           v_anchor="m", inset=G.INS_NONE, autofit="none")

    inject_footer(slide, sn)


def _rm(spPr):
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        spPr.remove(ln)
    etree.SubElement(spPr, qn("a:ln")).append(
        etree.fromstring(
            '<a:noFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
        )
    )
