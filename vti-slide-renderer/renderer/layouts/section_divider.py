# layouts/section_divider.py  ELayout G
from __future__ import annotations

from pptx.util import Emu
from lxml import etree
from pptx.oxml.ns import qn

from .. import geometry as G
from ..shapes import (
    add_textbox_styled, add_semantic_icon, add_rounded_rect,
    _set_gradient_fill_on_spPr, _set_solid_fill_on_spPr,
    add_decor_shape,
)
from ..footer import inject_footer
from ._common import remove_line


def render(prs, slide_data: dict):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    d = slide_data["design"]
    sn = slide_data["slide_number"]

    # ── Background ───────────────────────────────────────────────────
    bg = d.get("bg", {"from": "172759", "to": "1A3070"})
    bg_shape = slide.shapes.add_shape(
        1, Emu(0), Emu(0), Emu(G.SLIDE_W), Emu(G.SLIDE_H)
    )
    _set_gradient_fill_on_spPr(bg_shape._element.spPr,
                                bg["from"], bg["to"], bg.get("angle", 9000000))
    remove_line(bg_shape._element.spPr)
    bg_shape.text_frame.text = ""

    # ── Left bar ─────────────────────────────────────────────────────
    lb = d.get("left_bar", {"from": "2362B0", "to": "7FC236"})
    bar_w = G.pt(8)
    bar = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(bar_w), Emu(G.SLIDE_H))
    _set_gradient_fill_on_spPr(bar._element.spPr,
                                lb["from"], lb["to"], 18000000)
    remove_line(bar._element.spPr)
    bar.text_frame.text = ""

    # ── Decorative circle (top-right) ────────────────────────────────
    decor = d.get("decor", {"color": "1A3070", "opacity": 20})
    dec_sz = G.pt(300)
    dec_x = G.SLIDE_W - dec_sz // 2
    dec_y = -dec_sz // 4
    dec_shape = slide.shapes.add_shape(9, Emu(dec_x), Emu(dec_y),
                                        Emu(dec_sz), Emu(dec_sz))
    spPr = dec_shape._element.spPr
    opacity = decor.get("opacity", 20)
    alpha_emu = int((100 - opacity) / 100 * 100000)
    for tag in (qn("a:noFill"), qn("a:solidFill"), qn("a:gradFill")):
        for el in spPr.findall(tag):
            spPr.remove(el)
    solidFill = etree.SubElement(spPr, qn("a:solidFill"))
    srgb = etree.SubElement(solidFill, qn("a:srgbClr"), val=decor["color"].upper())
    etree.SubElement(srgb, qn("a:alpha"), val=str(alpha_emu))
    remove_line(spPr)
    dec_shape.text_frame.text = ""

    # ── Right card (tinted roundRect + top border + hero icon) ──────
    rc = d.get("right_card", {})
    card_w = G.pt(240)
    card_h = G.pt(360)
    card_x = G.SLIDE_W - card_w - G.pt(40)
    card_y = (G.SLIDE_H - card_h) // 2
    rc_bg = rc.get("bg", {"from": "EBF4FC", "to": "D6EAF8"})
    add_rounded_rect(slide, card_x, card_y, card_w, card_h,
                     rc_bg["from"], rc_bg["to"],
                     angle_emu=rc_bg.get("angle", 16200000),
                     border_hex=rc.get("border_top"),
                     border_w_pt=4)

    # Decor ellipse (Pattern 8) inside card
    add_decor_shape(slide,
                    card_x + G.pt(100), card_y + G.pt(220),
                    G.pt(180), G.pt(180),
                    rc.get("border_top", "4A9EE0"), alpha_percent=7)

    icon_def = rc.get("icon", {})
    if icon_def:
        icon_sz = icon_def.get("size_pt", 64)
        icon_x = card_x + (card_w - G.pt(icon_sz)) // 2
        icon_y = card_y + (card_h - G.pt(icon_sz)) // 2
        add_semantic_icon(slide, icon_x, icon_y, icon_sz,
                          icon_def.get("stroke", "4A9EE0"),
                          icon_def.get("type", "default"))

    # ── Text content ─────────────────────────────────────────────────
    tx = G.pt(28)
    text_w = card_x - G.pt(60)
    cy = G.SLIDE_H // 2 - G.pt(80)

    if d.get("section_number"):
        add_textbox_styled(slide, tx, cy, text_w, G.pt(30),
                           d["section_number"], bold=True, size_pt=22,
                           color_hex=d.get("section_number_color", "6A7FA0"),
                           v_anchor="m", inset=G.INS_NONE, autofit="none")
        cy += G.pt(36)

    if d.get("section_title"):
        add_textbox_styled(slide, tx, cy, text_w, G.pt(70),
                           d["section_title"], bold=True,
                           size_pt=G.FONT_SECTION_TITLE,
                           color_hex=d.get("section_title_color", "FFFFFF"),
                           v_anchor="t", inset=G.INS_NONE, autofit="norm", wrap=True)
        cy += G.pt(80)

    if d.get("subtitle"):
        add_textbox_styled(slide, tx, cy, text_w, G.pt(36),
                           d["subtitle"], size_pt=G.FONT_BODY,
                           color_hex=d.get("subtitle_color", "B8D0EE"),
                           v_anchor="m", inset=G.INS_NONE, autofit="none", wrap=True)

    inject_footer(slide, sn)
