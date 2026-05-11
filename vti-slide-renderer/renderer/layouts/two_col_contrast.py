# layouts/two_col_contrast.py — Layout D
from __future__ import annotations

from pptx.util import Emu
from .. import geometry as G
from ..shapes import add_textbox_styled, add_rounded_rect, add_dot_bullet, add_separator_line
from ..footer import inject_footer
from ._common import add_slide_bg, add_accent_bar, add_slide_header


def render(prs, slide_data: dict):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    d = slide_data["design"]
    sn = slide_data["slide_number"]

    add_slide_bg(slide, d.get("bg", {"from": "FFFFFF", "to": "EEF4FB"}))
    add_accent_bar(slide, d.get("accent_bar", {"from": "2362B0", "to": "7FC236"}))
    add_slide_header(slide, d.get("title", ""), d.get("breadcrumb", ""))

    col_w = G.card_w(2)
    pad = G.pt(16)

    for col_idx, (new_key, old_key) in enumerate(
            (("left_card", "left_col"), ("right_card", "right_col"))):
        col = d.get(new_key) or d.get(old_key, {})
        cx = G.card_x(col_idx, 2)
        cy = G.CONTENT_TOP
        ch = G.CONTENT_H
        bg = col.get("bg", {"from": "EBF4FC", "to": "D6EAF8"})
        add_rounded_rect(slide, cx, cy, col_w, ch,
                         bg["from"], bg["to"],
                         border_hex=col.get("border_top"), border_w_pt=3)

        inner_x = cx + pad
        inner_w = col_w - 2 * pad
        cur_y = cy + pad

        if col.get("header"):
            add_textbox_styled(slide, inner_x, cur_y, inner_w, G.pt(46),
                               col["header"], bold=True,
                               size_pt=G.FONT_CARD_HEADER,
                               color_hex=col.get("header_color", "172759"),
                               v_anchor="t", inset=G.INS_NONE,
                               autofit="norm", wrap=True)
            cur_y += G.pt(50)

            # Thin accent separator
            sep_color = col.get("border_top", "4A9EE0")
            add_separator_line(slide, inner_x, cur_y,
                               inner_w * 2 // 5, G.pt(1), sep_color)
            cur_y += G.pt(10)

        for bullet in col.get("bullets", []):
            dot_color = bullet.get("dot_color", "4A9EE0")
            add_dot_bullet(slide, inner_x, cur_y + G.pt(5), dot_color, 6)
            add_textbox_styled(slide, inner_x + G.pt(14), cur_y,
                               inner_w - G.pt(14), G.pt(30),
                               bullet.get("text", ""),
                               size_pt=G.FONT_BODY,
                               color_hex=col.get("text_color", "1C2D4F"),
                               v_anchor="t", inset=G.INS_NONE,
                               autofit="norm", wrap=True)
            cur_y += G.pt(32)

    inject_footer(slide, sn)
