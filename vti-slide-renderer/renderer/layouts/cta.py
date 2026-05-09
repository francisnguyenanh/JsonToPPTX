# layouts/cta.py  ELayout H
from __future__ import annotations

from pptx.util import Emu
from lxml import etree
from pptx.oxml.ns import qn

from .. import geometry as G
from ..shapes import (
    add_textbox_styled, add_badge, add_separator_line,
    _set_gradient_fill_on_spPr, _set_solid_fill_on_spPr,
    add_decor_shape,
)
from ..footer import inject_footer
from ._common import remove_line


def render(prs, slide_data: dict):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    d = slide_data["design"]
    sn = slide_data["slide_number"]

    # ── Background ───────────────────────────────────────────────────
    bg = d.get("bg", {"from": "2362B0", "to": "172759"})
    bg_shape = slide.shapes.add_shape(
        1, Emu(0), Emu(0), Emu(G.SLIDE_W), Emu(G.SLIDE_H)
    )
    _set_gradient_fill_on_spPr(bg_shape._element.spPr,
                                bg["from"], bg["to"], bg.get("angle", 16200000))
    remove_line(bg_shape._element.spPr)
    bg_shape.text_frame.text = ""

    # ── Decor circles (Pattern 8) — top-right and bottom-left ───────
    decor_color = d.get("decor_color", "2362B0")
    dec_sz = G.pt(300)
    # Top-right corner circle
    add_decor_shape(slide, G.SLIDE_W - dec_sz // 2, -dec_sz // 4,
                    dec_sz, dec_sz, decor_color, alpha_percent=7)
    # Bottom-left corner circle (smaller, accent colour)
    dec_sz2 = G.pt(220)
    add_decor_shape(slide, -dec_sz2 // 3, G.SLIDE_H - dec_sz2 * 3 // 4,
                    dec_sz2, dec_sz2, "7FC236", alpha_percent=6)

    # ── Heading ──────────────────────────────────────────────────────
    add_textbox_styled(slide, G.CONTENT_X, G.pt(55), G.CONTENT_W, G.pt(55),
                       d.get("heading", "Next Steps"), bold=True,
                       size_pt=G.FONT_CTA_HEADING,
                       color_hex=d.get("heading_color", "FFFFFF"),
                       v_anchor="m", inset=G.INS_NONE, autofit="none")

    # ── Divider ──────────────────────────────────────────────────────
    div_color = d.get("divider_color", "FFFFFF")
    add_separator_line(slide, G.CONTENT_X, G.pt(115), G.CONTENT_W, G.pt(1), div_color)

    # ── Items ────────────────────────────────────────────────────────
    items = d.get("items", [])
    n = len(items)
    item_w = G.card_w(n)

    for i, item in enumerate(items):
        ix = G.card_x(i, n)
        iy = G.CONTENT_TOP + G.pt(10)

        badge_color = item.get("badge_color", "7FC236")
        add_badge(slide, ix, iy, str(i + 1), badge_color, size_pt=18)

        title_y = iy + G.pt(26)
        add_textbox_styled(slide, ix, title_y, item_w, G.pt(24),
                           item.get("title", ""), bold=True, size_pt=G.FONT_HEADER,
                           color_hex=item.get("text_color", "FFFFFF"),
                           v_anchor="t", inset=G.INS_NONE, autofit="none", wrap=True)

        desc_y = title_y + G.pt(28)
        add_textbox_styled(slide, ix, desc_y, item_w,
                           G.CONTENT_H - G.pt(56),
                           item.get("description", ""),
                           size_pt=G.FONT_BODY,
                           color_hex=item.get("text_color", "FFFFFF"),
                           v_anchor="t", inset=G.INS_NONE, autofit="norm", wrap=True)

    inject_footer(slide, sn)
