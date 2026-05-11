# layouts/narrative.py — Layout F (Left-Text Right-Visual)
from __future__ import annotations

from pptx.util import Emu
from .. import geometry as G
from ..shapes import (
    add_textbox_styled, add_rounded_rect, add_dot_bullet,
    add_semantic_icon,
)
from ..footer import inject_footer
from ._common import add_slide_bg, add_accent_bar, add_slide_header


def render(prs, slide_data: dict):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    d = slide_data["design"]
    sn = slide_data["slide_number"]

    add_slide_bg(slide, d.get("bg", {"from": "FFFFFF", "to": "EEF4FB"}))
    add_accent_bar(slide, d.get("accent_bar", {"from": "2362B0", "to": "7FC236"}))
    add_slide_header(slide, d.get("title", ""), d.get("breadcrumb", ""),
                     title_color=d.get("title_color", "172759"),
                     breadcrumb_color=d.get("breadcrumb_color", "6A7FA0"))

    # Zone split: 57% left panel, 5% gap, 38% right panel
    right_w = G.pt(340)
    gap = G.pt(16)
    left_w = G.CONTENT_W - right_w - gap

    left_x = G.CONTENT_X
    right_x = left_x + left_w + gap

    # ── Left panel: body text + optional bullets ──────────────────────
    left_panel = d.get("left_panel", {})
    body_text = left_panel.get("body_text", d.get("body", ""))
    text_color = left_panel.get("text_color", d.get("text_color", "1C2D4F"))

    cur_y = G.CONTENT_TOP
    bullets = left_panel.get("bullets", [])

    if body_text:
        # Reserve space for bullets: 32pt per bullet + 8pt gap before first bullet
        bullet_reserve = len(bullets) * G.pt(32) + (G.pt(12) if bullets else 0)
        body_h = max(G.CONTENT_H - bullet_reserve, G.pt(80))
        add_textbox_styled(
            slide, left_x, cur_y, left_w, body_h,
            body_text,
            size_pt=G.FONT_BODY, color_hex=text_color,
            v_anchor="t", inset=G.INS_CARD,
            autofit="norm", wrap=True,
            line_spacing=1.5
        )
        cur_y += body_h + G.pt(10)

    for bullet in bullets:
        add_dot_bullet(slide, left_x + G.pt(4), cur_y + G.pt(6),
                       bullet.get("dot_color", "4A9EE0"), 6)
        add_textbox_styled(slide, left_x + G.pt(18), cur_y,
                           left_w - G.pt(18), G.pt(28),
                           bullet.get("text", ""),
                           size_pt=G.FONT_BODY, color_hex=text_color,
                           v_anchor="t", inset=G.INS_NONE,
                           autofit="norm", wrap=True)
        cur_y += G.pt(32)

    # ── Right panel: tinted card + icon + optional callout ────────────
    right_panel = d.get("right_panel", {})
    if right_panel:
        rp_bg = right_panel.get("bg", {"from": "EBF4FC", "to": "D6EAF8"})
        add_rounded_rect(slide, right_x, G.CONTENT_TOP, right_w, G.CONTENT_H,
                         rp_bg["from"], rp_bg["to"],
                         angle_emu=rp_bg.get("angle", 16200000),
                         border_hex=right_panel.get("border_top"),
                         border_w_pt=3)

        icon_def = right_panel.get("icon", {})
        callout_number = right_panel.get("callout_number", "")
        callout_label = right_panel.get("callout_label", "")

        if callout_number:
            add_textbox_styled(
                slide,
                right_x + G.pt(8), G.CONTENT_TOP + G.pt(60),
                right_w - G.pt(16), G.pt(80),
                callout_number,
                bold=True, size_pt=G.FONT_STAT_BIG,
                color_hex=right_panel.get("border_top", "2362B0"),
                align="center", v_anchor="m", inset=G.INS_NONE, autofit="norm"
            )
            if callout_label:
                add_textbox_styled(
                    slide,
                    right_x + G.pt(8), G.CONTENT_TOP + G.pt(145),
                    right_w - G.pt(16), G.pt(28),
                    callout_label,
                    size_pt=G.FONT_STAT_LABEL,
                    color_hex="6A7FA0",
                    align="center", v_anchor="m", inset=G.INS_NONE, autofit="norm"
                )
        elif icon_def:
            icon_sz = icon_def.get("size_pt", 56)
            icon_x = right_x + (right_w - G.pt(icon_sz)) // 2
            icon_y = G.CONTENT_TOP + (G.CONTENT_H - G.pt(icon_sz)) // 2
            add_semantic_icon(slide, icon_x, icon_y, icon_sz,
                              icon_def.get("stroke", "4A9EE0"),
                              icon_def.get("type", "default"))

    inject_footer(slide, sn)
