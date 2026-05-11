# layouts/toc.py — Layout J (Table of Contents / Agenda)
from __future__ import annotations

from pptx.util import Emu
from .. import geometry as G
from ..shapes import add_textbox_styled, add_badge, add_separator_line
from ..footer import inject_footer
from ._common import add_slide_bg, add_accent_bar, add_slide_header


def render(prs, slide_data: dict):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    d = slide_data["design"]
    sn = slide_data["slide_number"]

    add_slide_bg(slide, d.get("bg", {"from": "FFFFFF", "to": "EEF4FB"}))
    add_accent_bar(slide, d.get("accent_bar", {"from": "2362B0", "to": "7FC236"}))
    add_slide_header(slide, d.get("title", "Agenda"), d.get("breadcrumb", ""),
                     title_color=d.get("title_color", "172759"))

    sep_color = d.get("separator_color", "B8D0EE")

    # Detect format: new (left_items/right_items) vs legacy (items)
    left_items = d.get("left_items", [])
    right_items = d.get("right_items", [])

    if left_items or right_items:
        # ── New format: left_items + right_items columns ──────────────
        col_w = G.card_w(2)
        left_x = G.card_x(0, 2)
        right_x = G.card_x(1, 2)

        # Vertical separator
        add_separator_line(slide,
                           left_x + col_w + G.pt(2), G.CONTENT_TOP,
                           G.pt(2), G.CONTENT_H, sep_color)

        def _draw_col(items, start_x, start_idx):
            row_h = G.pt(80)
            for i, item in enumerate(items):
                iy = G.CONTENT_TOP + i * row_h
                badge_color = item.get("badge_color", "2362B0")
                badge_num = str(start_idx + i + 1)
                add_badge(slide, start_x, iy + G.pt(11), badge_num, badge_color, size_pt=21)

                text_x = start_x + G.pt(29)
                text_w = col_w - G.pt(32)

                # Section title
                title_color = item.get("title_color", "172759")
                add_textbox_styled(
                    slide, text_x, iy + G.pt(8), text_w, G.pt(37),
                    item.get("section_title", item.get("label", "")),
                    bold=True, size_pt=G.FONT_CARD_HEADER,
                    color_hex=title_color,
                    v_anchor="m", inset=G.INS_NONE, autofit="none"
                )

                # Slide range
                slide_range = item.get("slide_range", item.get("description", ""))
                if slide_range:
                    range_color = item.get("range_color", item.get("sub_color", "6A7FA0"))
                    add_textbox_styled(
                        slide, text_x, iy + G.pt(45), text_w, G.pt(27),
                        slide_range,
                        size_pt=G.FONT_SMALL, color_hex=range_color,
                        v_anchor="t", inset=G.INS_NONE, autofit="none"
                    )

        _draw_col(left_items, left_x, 0)
        _draw_col(right_items, right_x, len(left_items))

    else:
        # ── Legacy format: flat items list ────────────────────────────
        items = d.get("items", [])
        n = len(items)
        cols = 2 if n > 4 else 1
        col_w = G.card_w(cols)

        for i, item in enumerate(items):
            col = i % cols
            row = i // cols
            ix = G.card_x(col, cols)
            iy = G.CONTENT_TOP + row * G.pt(59)

            badge_color = item.get("badge_color", "2362B0")
            add_badge(slide, ix, iy, str(i + 1), badge_color, size_pt=21)

            add_textbox_styled(
                slide, ix + G.pt(29), iy, col_w - G.pt(32), G.pt(40),
                item.get("label", ""), bold=True, size_pt=G.FONT_CARD_HEADER,
                color_hex=item.get("text_color", "172759"),
                v_anchor="m", inset=G.INS_NONE, autofit="none"
            )

            if item.get("description"):
                add_textbox_styled(
                    slide, ix + G.pt(29), iy + G.pt(29),
                    col_w - G.pt(32), G.pt(24),
                    item["description"], size_pt=G.FONT_SMALL,
                    color_hex=item.get("sub_color", "6A7FA0"),
                    v_anchor="t", inset=G.INS_NONE, autofit="none"
                )

    inject_footer(slide, sn)
