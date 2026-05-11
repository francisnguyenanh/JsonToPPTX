# shapes.py — Core shape builders using python-pptx + lxml OOXML injection
from __future__ import annotations

from lxml import etree
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn, nsmap

from . import geometry as G


# ── Helper: parse hex color ───────────────────────────────────────────

def _rgb(hex6: str) -> RGBColor:
    h = hex6.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _hex_upper(hex6: str) -> str:
    return hex6.lstrip("#").upper()


# ── Gradient fill via OOXML ───────────────────────────────────────────

def add_gradient_fill(shape, from_hex: str, to_hex: str, angle_emu: int, stops=None):
    """
    Apply gradient fill to any shape via OOXML injection.
    stops: optional list of (position_0_to_100000, hex) for multi-stop gradients.
    """
    spPr = shape.fill._xPr  # sp element for shapes; works for shapes
    # Remove any existing fill
    for tag in (qn("a:noFill"), qn("a:solidFill"), qn("a:gradFill"),
                qn("a:pattFill"), qn("a:blipFill")):
        for el in spPr.findall(tag):
            spPr.remove(el)

    gradFill = etree.SubElement(spPr, qn("a:gradFill"))
    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))

    if stops is None:
        stops = [(0, from_hex), (100000, to_hex)]

    for pos, hex_color in stops:
        gs = etree.SubElement(gsLst, qn("a:gs"), pos=str(pos))
        srgbClr = etree.SubElement(gs, qn("a:srgbClr"), val=_hex_upper(hex_color))

    lin = etree.SubElement(gradFill, qn("a:lin"), ang=str(angle_emu), scaled="0")
    return gradFill


def _set_gradient_fill_on_spPr(spPr, from_hex: str, to_hex: str, angle_emu: int, stops=None):
    """Apply gradient fill directly on an spPr element (for shapes without .fill)."""
    for tag in (qn("a:noFill"), qn("a:solidFill"), qn("a:gradFill"),
                qn("a:pattFill"), qn("a:blipFill")):
        for el in spPr.findall(tag):
            spPr.remove(el)

    gradFill = etree.SubElement(spPr, qn("a:gradFill"))
    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))

    if stops is None:
        stops = [(0, from_hex), (100000, to_hex)]

    for pos, hex_color in stops:
        gs = etree.SubElement(gsLst, qn("a:gs"), pos=str(pos))
        etree.SubElement(gs, qn("a:srgbClr"), val=_hex_upper(hex_color))

    etree.SubElement(gradFill, qn("a:lin"), ang=str(angle_emu), scaled="0")


def _set_solid_fill_on_spPr(spPr, hex_color: str, alpha: int = None):
    """Apply solid fill on spPr element. alpha: 0-100000 (100000=opaque)."""
    for tag in (qn("a:noFill"), qn("a:solidFill"), qn("a:gradFill"),
                qn("a:pattFill"), qn("a:blipFill")):
        for el in spPr.findall(tag):
            spPr.remove(el)
    solidFill = etree.SubElement(spPr, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"), val=_hex_upper(hex_color))
    if alpha is not None:
        etree.SubElement(srgbClr, qn("a:alpha"), val=str(alpha))


# ── Rounded rectangle ────────────────────────────────────────────────

def add_rounded_rect(slide, x: int, y: int, w: int, h: int,
                     from_hex: str, to_hex: str,
                     angle_emu: int = 16200000,
                     border_hex: str = None, border_w_pt: float = 0):
    """
    Add roundRect shape (adj=20000, ~16% corner radius).
    Optional top border as separate thin rect if border_hex provided.
    Returns the shape object.
    """
    from pptx.util import Emu as _Emu
    shape = slide.shapes.add_shape(
        1,
        _Emu(x), _Emu(y), _Emu(w), _Emu(h)
    )

    # Set preset geometry to roundRect with a slightly more rounded corner
    spPr = shape._element.spPr
    prstGeom = spPr.find(qn("a:prstGeom"))
    if prstGeom is not None:
        prstGeom.set("prst", "roundRect")
        avLst = prstGeom.find(qn("a:avLst"))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn("a:avLst"))
        else:
            for gd in avLst.findall(qn("a:gd")):
                avLst.remove(gd)
        etree.SubElement(avLst, qn("a:gd"), name="adj", fmla="val 20000")

    # Remove line (border)
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        spPr.remove(ln)
    ln_el = etree.SubElement(spPr, qn("a:ln"))
    etree.SubElement(ln_el, qn("a:noFill"))

    # Gradient fill
    _set_gradient_fill_on_spPr(spPr, from_hex, to_hex, angle_emu)

    # Optional top border line
    if border_hex:
        bw = G.pt(border_w_pt) if border_w_pt else G.pt(3)
        border = slide.shapes.add_shape(1, _Emu(x), _Emu(y), _Emu(w), _Emu(bw))
        b_spPr = border._element.spPr
        _set_solid_fill_on_spPr(b_spPr, border_hex)
        b_ln = b_spPr.find(qn("a:ln"))
        if b_ln is not None:
            b_spPr.remove(b_ln)
        etree.SubElement(b_spPr, qn("a:ln")).append(
            etree.fromstring('<a:noFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
        )
        prstGeomB = b_spPr.find(qn("a:prstGeom"))
        if prstGeomB is not None:
            prstGeomB.set("prst", "rect")

    return shape


# ── Styled textbox ────────────────────────────────────────────────────

def add_textbox_styled(slide, x: int, y: int, w: int, h: int,
                       text: str,
                       font_name: str = G.FONT_NAME,
                       size_pt: float = G.FONT_BODY,
                       color_hex: str = "000000",
                       bold: bool = False,
                       italic: bool = False,
                       align: str = "left",
                       v_anchor: str = "t",
                       wrap: bool = True,
                       inset: dict = None,
                       autofit: str = "norm",
                       line_spacing: float = None,
                       lang: str = "vi-VN"):
    """
    Add textbox with full bodyPr set explicitly.
    Never leaves <a:bodyPr/> empty.
    """
    from pptx.util import Emu as _Emu
    txBox = slide.shapes.add_textbox(_Emu(x), _Emu(y), _Emu(w), _Emu(h))
    tf = txBox.text_frame

    # Build bodyPr attributes
    ins = inset if inset else G.INS_NONE
    wrap_val = "square" if wrap else "none"
    anchor_map = {"t": "t", "m": "ctr", "b": "b", "ctr": "ctr"}
    anchor_val = anchor_map.get(v_anchor, "t")

    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    bodyPr.set("wrap", wrap_val)
    bodyPr.set("lIns", str(ins.get("lIns", 0)))
    bodyPr.set("rIns", str(ins.get("rIns", 0)))
    bodyPr.set("tIns", str(ins.get("tIns", 0)))
    bodyPr.set("bIns", str(ins.get("bIns", 0)))
    bodyPr.set("anchor", anchor_val)
    bodyPr.set("anchorCtr", "0")

    # Remove existing autofit children
    for tag in (qn("a:normAutofit"), qn("a:noAutofit"), qn("a:spAutoFit")):
        for el in bodyPr.findall(tag):
            bodyPr.remove(el)

    if autofit == "norm":
        etree.SubElement(bodyPr, qn("a:normAutofit"))
    elif autofit == "sp":
        etree.SubElement(bodyPr, qn("a:spAutoFit"))
    else:
        etree.SubElement(bodyPr, qn("a:noAutofit"))

    # Paragraph
    tf.word_wrap = wrap
    p = tf.paragraphs[0]

    # Alignment
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                 "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
    p.alignment = align_map.get(align, PP_ALIGN.LEFT)

    if line_spacing:
        from pptx.util import Pt as _Pt
        from pptx.oxml.ns import qn as _qn
        pPr = p._pPr
        if pPr is None:
            pPr = etree.SubElement(p._p, qn("a:pPr"))
        lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
        spcPts = etree.SubElement(lnSpc, qn("a:spcPts"))
        spcPts.set("val", str(int(line_spacing * 100)))

    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.color.rgb = _rgb(color_hex)
    run.font.bold = bold
    run.font.italic = italic

    # Language
    rPr = run._r.find(qn("a:rPr"))
    if rPr is None:
        rPr = run._r.find(qn("a:rPr"))
    if rPr is not None:
        rPr.set("lang", lang)

    return txBox


# ── Badge (circular ellipse) ──────────────────────────────────────────

def add_badge(slide, x: int, y: int, number_or_text: str,
              bg_hex: str, text_hex: str = "FFFFFF", size_pt: float = 40):
    """Circular badge (ellipse) with centered number/text."""
    from pptx.util import Emu as _Emu
    sz = G.pt(size_pt)
    shape = slide.shapes.add_shape(9, _Emu(x), _Emu(y), _Emu(sz), _Emu(sz))  # 9 = oval

    spPr = shape._element.spPr
    _set_solid_fill_on_spPr(spPr, bg_hex)
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        spPr.remove(ln)
    etree.SubElement(spPr, qn("a:ln")).append(
        etree.fromstring('<a:noFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
    )

    tf = shape.text_frame
    tf.word_wrap = False
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    bodyPr.set("anchor", "ctr")
    bodyPr.set("anchorCtr", "1")
    bodyPr.set("lIns", "0")
    bodyPr.set("rIns", "0")
    bodyPr.set("tIns", "0")
    bodyPr.set("bIns", "0")

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number_or_text)
    run.font.name = G.FONT_NAME
    run.font.size = Pt(size_pt * 0.4)
    run.font.color.rgb = _rgb(text_hex)
    run.font.bold = True

    return shape


# ── Dot bullet ────────────────────────────────────────────────────────

def add_dot_bullet(slide, x: int, y: int, dot_hex: str, size_pt: float = 6):
    """Small filled circle for bullet points."""
    from pptx.util import Emu as _Emu
    sz = G.pt(size_pt)
    shape = slide.shapes.add_shape(9, _Emu(x), _Emu(y), _Emu(sz), _Emu(sz))
    spPr = shape._element.spPr
    _set_solid_fill_on_spPr(spPr, dot_hex)
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        spPr.remove(ln)
    etree.SubElement(spPr, qn("a:ln")).append(
        etree.fromstring('<a:noFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
    )
    shape.text_frame.text = ""
    return shape


# ── Separator line ────────────────────────────────────────────────────

def add_separator_line(slide, x: int, y: int, w: int, h: int, color_hex: str):
    """Thin rect used as horizontal or vertical separator."""
    from pptx.util import Emu as _Emu
    shape = slide.shapes.add_shape(1, _Emu(x), _Emu(y), _Emu(w), _Emu(h))
    spPr = shape._element.spPr
    _set_solid_fill_on_spPr(spPr, color_hex)
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        spPr.remove(ln)
    etree.SubElement(spPr, qn("a:ln")).append(
        etree.fromstring('<a:noFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
    )
    shape.text_frame.text = ""
    return shape


# ── Semantic icon (filled OOXML presets + glow background) ────────────

def add_semantic_icon(slide, x: int, y: int, size_pt: float,
                      stroke_hex: str, icon_type: str):
    """
    Professional icon: semi-transparent glow circle + filled preset shape symbol.
    No bitmaps, no PIL — pure OOXML geometry.
    Returns list of shapes added.
    """
    from pptx.util import Emu as _Emu

    sz = G.pt(size_pt)
    result = []
    _NOFILL = '<a:noFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'

    def _prep(s, preset):
        spPr = s._element.spPr
        pg = spPr.find(qn("a:prstGeom"))
        if pg is not None:
            pg.set("prst", preset)
            av = pg.find(qn("a:avLst"))
            if av is None:
                etree.SubElement(pg, qn("a:avLst"))
            else:
                for gd in av.findall(qn("a:gd")):
                    av.remove(gd)
        for tag in (qn("a:noFill"), qn("a:solidFill"), qn("a:gradFill")):
            for el in spPr.findall(tag):
                spPr.remove(el)
        old_ln = spPr.find(qn("a:ln"))
        if old_ln is not None:
            spPr.remove(old_ln)
        return spPr

    def filled(preset, sx, sy, sw, sh, fill_hex, alpha=None):
        s = slide.shapes.add_shape(1, _Emu(sx), _Emu(sy), _Emu(sw), _Emu(sh))
        spPr = _prep(s, preset)
        sf = etree.SubElement(spPr, qn("a:solidFill"))
        srgb = etree.SubElement(sf, qn("a:srgbClr"), val=_hex_upper(fill_hex))
        if alpha is not None:
            etree.SubElement(srgb, qn("a:alpha"), val=str(alpha))
        ln = etree.SubElement(spPr, qn("a:ln"))
        etree.SubElement(ln, qn("a:noFill"))
        s.text_frame.text = ""
        result.append(s)
        return s

    def stroked(preset, sx, sy, sw, sh, color_hex, w_pt=2.5):
        s = slide.shapes.add_shape(1, _Emu(sx), _Emu(sy), _Emu(sw), _Emu(sh))
        spPr = _prep(s, preset)
        etree.SubElement(spPr, qn("a:noFill"))
        ln = etree.SubElement(spPr, qn("a:ln"), w=str(G.pt(w_pt)))
        sf = etree.SubElement(ln, qn("a:solidFill"))
        etree.SubElement(sf, qn("a:srgbClr"), val=_hex_upper(color_hex))
        s.text_frame.text = ""
        result.append(s)
        return s

    # Background glow circle (semi-transparent)
    filled("ellipse", x, y, sz, sz, stroke_hex, alpha=18000)

    # Inner icon zone with padding
    pad = sz // 5
    ix, iy, isz = x + pad, y + pad, sz - 2 * pad

    if icon_type == "ai":
        # Hexagon (tech/AI feel) + white cut-out center
        filled("hexagon", ix, iy, isz, isz, stroke_hex)
        q = isz // 3
        filled("ellipse", ix + q, iy + q, isz - 2 * q, isz - 2 * q, "FFFFFF")

    elif icon_type == "data":
        # Filled bar chart — 3 bars at varied heights
        bw = max((isz - G.pt(4)) // 3, G.pt(6))
        for i, hf in enumerate([0.55, 0.88, 0.40]):
            bh = max(int(isz * hf), G.pt(4))
            bx = ix + i * (bw + G.pt(2))
            by = iy + isz - bh
            filled("rect", bx, by, bw, bh, stroke_hex)

    elif icon_type == "settings":
        # 12-point star (gear-like) + white center
        filled("star12", ix, iy, isz, isz, stroke_hex)
        q = isz // 3
        filled("ellipse", ix + q, iy + q, isz - 2 * q, isz - 2 * q, "FFFFFF")

    elif icon_type == "flow":
        # Filled right-arrow
        ah = max(isz * 3 // 5, G.pt(8))
        filled("rightArrow", ix, iy + (isz - ah) // 2, isz, ah, stroke_hex)

    elif icon_type == "doc":
        # Rounded document + 3 white horizontal lines
        filled("roundRect", ix, iy, isz, isz, stroke_hex)
        lh = max(int(G.pt(1.5)), 1)
        lx = ix + isz // 6
        lw = isz * 2 // 3
        for j in range(3):
            ly = iy + isz // 5 + j * (isz // 5)
            filled("rect", lx, ly, lw, lh, "FFFFFF")

    elif icon_type == "team":
        # Two person silhouettes
        hs = isz * 2 // 5
        bh = max(isz - hs - G.pt(3), G.pt(4))
        filled("ellipse", ix, iy, hs, hs, stroke_hex)
        filled("roundRect", ix - G.pt(2), iy + hs + G.pt(3), hs + G.pt(4), bh, stroke_hex)
        p2x = ix + isz - hs
        filled("ellipse", p2x, iy, hs, hs, stroke_hex)
        filled("roundRect", p2x - G.pt(2), iy + hs + G.pt(3), hs + G.pt(4), bh, stroke_hex)

    elif icon_type == "check":
        # Filled circle + white checkmark lines
        filled("ellipse", ix, iy, isz, isz, stroke_hex)
        lh = max(int(G.pt(2)), 1)
        filled("rect", ix + isz // 5, iy + isz * 3 // 5, isz // 4, lh, "FFFFFF")
        filled("rect", ix + isz * 2 // 5, iy + isz // 3, lh, isz // 3, "FFFFFF")

    elif icon_type == "chart":
        # Rounded background + white bar chart
        filled("roundRect", ix, iy, isz, isz, stroke_hex)
        bw = max(isz // 7, G.pt(4))
        for j, hf in enumerate([0.40, 0.72, 0.55, 0.88]):
            bh = int(isz * hf)
            bx = ix + G.pt(3) + j * (bw + G.pt(2))
            by = iy + isz - bh - G.pt(3)
            if bx + bw <= ix + isz - G.pt(2):
                filled("rect", bx, by, bw, bh, "FFFFFF")

    elif icon_type == "lock":
        # Lock body + stroked shackle arc
        body_h = isz * 2 // 3
        filled("roundRect", ix, iy + isz - body_h, isz, body_h, stroke_hex)
        shw = isz * 5 // 8
        stroked("ellipse", ix + (isz - shw) // 2, iy, shw, isz * 2 // 3, stroke_hex, 2.5)

    elif icon_type == "star":
        filled("star5", ix, iy, isz, isz, stroke_hex)

    elif icon_type == "lightning":
        filled("lightningBolt", ix, iy, isz, isz, stroke_hex)

    elif icon_type == "arrow_up":
        filled("upArrow", ix, iy, isz, isz, stroke_hex)

    elif icon_type == "pentagon":
        filled("pentagon", ix, iy, isz, isz, stroke_hex)

    else:
        # default — filled diamond
        filled("diamond", ix, iy, isz, isz, stroke_hex)

    return result


# ── Decorative background shape (Pattern 8) ──────────────────────────

def add_decor_shape(slide, x: int, y: int, w: int, h: int,
                    color_hex: str, alpha_percent: int = 8):
    """
    Large faint ellipse for visual depth (Pattern 8).
    alpha_percent: 5-15. Append BEFORE content shapes (z-order).
    """
    from pptx.util import Emu as _Emu
    shape = slide.shapes.add_shape(9, _Emu(x), _Emu(y), _Emu(w), _Emu(h))
    spPr = shape._element.spPr
    alpha_emu = int(alpha_percent * 1000)
    for tag in (qn("a:noFill"), qn("a:solidFill"), qn("a:gradFill")):
        for el in spPr.findall(tag):
            spPr.remove(el)
    solidFill = etree.SubElement(spPr, qn("a:solidFill"))
    srgb = etree.SubElement(solidFill, qn("a:srgbClr"), val=_hex_upper(color_hex))
    etree.SubElement(srgb, qn("a:alpha"), val=str(alpha_emu))
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        spPr.remove(ln)
    etree.SubElement(spPr, qn("a:ln")).append(
        etree.fromstring('<a:noFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
    )
    shape.text_frame.text = ""
    return shape
