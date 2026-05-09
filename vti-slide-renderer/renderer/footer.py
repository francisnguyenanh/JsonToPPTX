# footer.py  EFooter injector, mandatory on every slide
from . import geometry as G
from .shapes import add_textbox_styled, _set_solid_fill_on_spPr
from lxml import etree
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def inject_footer(slide, slide_number: int):
    """
    Insert footer shapes into slide. Called LAST after all content shapes.
    Fixed positions from geometry.py constants  Ecoordinates never modified.

    Shape 1: Filled rect (slide number badge)
             fill = #408DD6
             text = slide_number, Arial 18pt bold white, centered

    Shape 2: Copyright textbox (transparent)
             text = "Copyright © 2026 VTI All rights reserved"
             Arial 10pt, color #808080, wrap=none
    """
    # Shape 1  ESlide number badge
    badge = slide.shapes.add_shape(
        1,
        Emu(G.FOOTER_NUM_X), Emu(G.FOOTER_NUM_Y),
        Emu(G.FOOTER_NUM_W), Emu(G.FOOTER_NUM_H)
    )
    spPr = badge._element.spPr
    _set_solid_fill_on_spPr(spPr, "408DD6")
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        spPr.remove(ln)
    etree.SubElement(spPr, qn("a:ln")).append(
        etree.fromstring('<a:noFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
    )

    tf = badge.text_frame
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    bodyPr.set("anchor", "ctr")
    bodyPr.set("anchorCtr", "1")
    bodyPr.set("lIns", str(G.INS_FOOTER["lIns"]))
    bodyPr.set("rIns", str(G.INS_FOOTER["rIns"]))
    bodyPr.set("tIns", str(G.INS_FOOTER["tIns"]))
    bodyPr.set("bIns", str(G.INS_FOOTER["bIns"]))
    bodyPr.set("wrap", "none")
    for tag in (qn("a:normAutofit"), qn("a:noAutofit"), qn("a:spAutoFit")):
        for el in bodyPr.findall(tag):
            bodyPr.remove(el)
    etree.SubElement(bodyPr, qn("a:noAutofit"))

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(slide_number)
    run.font.name = G.FONT_NAME
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Shape 2  ECopyright text
    copy_box = slide.shapes.add_textbox(
        Emu(G.FOOTER_COPY_X), Emu(G.FOOTER_COPY_Y),
        Emu(G.FOOTER_COPY_W), Emu(G.FOOTER_COPY_H)
    )
    ctf = copy_box.text_frame
    c_bodyPr = ctf._txBody.find(qn("a:bodyPr"))
    c_bodyPr.set("wrap", "none")
    c_bodyPr.set("anchor", "ctr")
    c_bodyPr.set("anchorCtr", "0")
    c_bodyPr.set("lIns", str(G.INS_FOOTER["lIns"]))
    c_bodyPr.set("rIns", str(G.INS_FOOTER["rIns"]))
    c_bodyPr.set("tIns", "0")
    c_bodyPr.set("bIns", "0")
    for tag in (qn("a:normAutofit"), qn("a:noAutofit"), qn("a:spAutoFit")):
        for el in c_bodyPr.findall(tag):
            c_bodyPr.remove(el)
    etree.SubElement(c_bodyPr, qn("a:noAutofit"))

    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.LEFT
    crun = cp.add_run()
    crun.text = "Copyright \u00a9 2026 VTI All rights reserved"
    crun.font.name = G.FONT_NAME
    crun.font.size = Pt(10)
    crun.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
