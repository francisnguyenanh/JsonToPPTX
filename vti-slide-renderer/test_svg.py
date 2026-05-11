"""Quick SVG parser smoke-test."""
from renderer.svg_parser import parse_svg_to_slide
from pptx import Presentation
from pptx.util import Emu
from renderer import geometry as G
import io, pathlib

prs = Presentation()
prs.slide_width  = Emu(G.SLIDE_W)
prs.slide_height = Emu(G.SLIDE_H)

SVG1 = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">
  <defs>
    <linearGradient id="bg" x1="0" y1="0" x2="0" y2="1">
      <stop offset="0%" style="stop-color:#172759"/>
      <stop offset="100%" style="stop-color:#2362B0"/>
    </linearGradient>
  </defs>
  <rect width="1280" height="720" fill="url(#bg)"/>
  <rect x="0" y="0" width="14" height="720" fill="#2362B0"/>
  <rect x="570" y="80" width="410" height="600" rx="12" fill="#EBF4FC" stroke="#7FC236" stroke-width="4"/>
  <text x="51" y="150" font-size="20" font-weight="bold" font-family="Calibri" fill="#7FC236">VTI INNOVATION</text>
  <text x="51" y="270" font-size="64" font-weight="bold" font-family="Calibri" fill="white">AI Slide Generator</text>
  <circle cx="775" cy="380" r="80" fill="#2362B0" opacity="0.15"/>
  <line x1="51" y1="390" x2="407" y2="390" stroke="#2362B0" stroke-width="4"/>
</svg>"""

SVG2 = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">
  <rect width="1280" height="720" fill="#F5F6F8"/>
  <text x="640" y="180" font-size="48" font-weight="bold" font-family="Calibri" fill="#172759" text-anchor="middle">Content Slide</text>
  <g transform="translate(0, 40)">
    <rect x="100" y="240" width="300" height="180" rx="10" fill="#EBF4FC" stroke="#2362B0" stroke-width="2"/>
    <text x="250" y="325" font-size="22" font-family="Calibri" fill="#172759" text-anchor="middle">Card 1</text>
    <rect x="490" y="240" width="300" height="180" rx="10" fill="#E3F6F0" stroke="#1E9E8A" stroke-width="2"/>
    <text x="640" y="325" font-size="22" font-family="Calibri" fill="#172759" text-anchor="middle">Card 2</text>
    <rect x="880" y="240" width="300" height="180" rx="10" fill="#FEF5E0" stroke="#E8A020" stroke-width="2"/>
    <text x="1030" y="325" font-size="22" font-family="Calibri" fill="#172759" text-anchor="middle">Card 3</text>
  </g>
</svg>"""

parse_svg_to_slide(prs, SVG1)
parse_svg_to_slide(prs, SVG2)

buf = io.BytesIO()
prs.save(buf)
pptx_bytes = buf.getvalue()

pathlib.Path("output/svg_test.pptx").write_bytes(pptx_bytes)
print(f"OK — output/svg_test.pptx ({len(pptx_bytes)//1024} KB, 2 slides)")
print(f"Slide 1 shapes: {len(prs.slides[0].shapes)}")
print(f"Slide 2 shapes: {len(prs.slides[1].shapes)}")
